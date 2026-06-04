"""
generate_eda_ppt.py
-------------------
Generates  outputs/FinnUp_EDA_Report.pptx
  – a fully self-contained EDA summary deck from FinnUp_Borrowers.xlsx

Run:
    python generate_eda_ppt.py
"""

from __future__ import annotations

import io
import os
import sys
import warnings
from pathlib import Path

import matplotlib
matplotlib.use("Agg")          # headless – no GUI window
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
import seaborn as sns
warnings.filterwarnings("ignore")

# ── pptx imports ──────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))
from src.data.loader import (
    load_all, SHEET, summary,
    get_profile_with_bank, get_kpi_pivot, get_director_summary,
)

# ── Config ────────────────────────────────────────────────────────────────────
EXCEL_PATH   = r"C:\Users\Ganesh.Bisht\Downloads\FinnUp_Borrowers.xlsx"
OUTPUT_PATH  = Path("outputs") / "FinnUp_EDA_Report.pptx"
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# Palette
NAVY    = RGBColor(0x1B, 0x3A, 0x6B)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
SALMON  = RGBColor(0xEF, 0x44, 0x44)

sns.set_theme(style="whitegrid", palette="muted", font_scale=1.0)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

# ── Helper: fig → BytesIO PNG ─────────────────────────────────────────────────
def fig_to_bytes(fig, dpi: int = 130) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

# ── Helper: add a slide with navy header bar ──────────────────────────────────
def add_slide(prs: Presentation, title: str,
              subtitle: str = "") -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Top header bar
    bar = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SW, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    tf.margin_left  = Inches(0.35)
    tf.margin_top   = Inches(0.18)

    # Subtitle (smaller, teal)
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.35), Inches(0.82), Inches(12), Pt(18)
        )
        sub_tf = sub_box.text_frame
        sub_p  = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.size = Pt(11)
        sub_run.font.color.rgb = TEAL
        sub_run.font.italic = True

    # Thin accent line below header
    line = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), SW, Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    return slide


def add_textbox(slide, left, top, width, height, text,
                bold=False, size=11, color=DARK,
                align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text  = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, left, top, width, height,
                   items: list, size=11, title=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.bold  = True
        r.font.size  = Pt(size + 1)
        r.font.color.rgb = NAVY
    for item in items:
        p = tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size  = Pt(size)
        r.font.color.rgb = DARK
    return tb


def add_stat_box(slide, left, top, width, height,
                 value: str, label: str,
                 bg: RGBColor = NAVY):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.12)
    # value
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.bold  = True
    r.font.size  = Pt(26)
    r.font.color.rgb = WHITE
    # label
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size  = Pt(10)
    r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)


# ═══════════════════════════════════════════════════════════════════════════════
# LOAD DATA
# ═══════════════════════════════════════════════════════════════════════════════
print("Loading data …")
data = load_all(EXCEL_PATH, verbose=True)
prof = data[SHEET.PROFILE]
dirs = data[SHEET.DIRECTORS]
bank = data[SHEET.BANK]
loans = data[SHEET.LOANS]
key_kpis = ["DSCR (Avg/Min)", "Current Ratio", "Debt Equity Ratio",
            "Net Profit Margin (%)", "PBDITA Margin (%)", "Debt/PBITDA Ratio"]
kpi_pivot   = get_kpi_pivot(data, kpi_labels=key_kpis, latest_year_only=True)
dir_sum     = get_director_summary(data)
master      = get_profile_with_bank(data)
if not kpi_pivot.empty:
    master = master.merge(kpi_pivot, on="Borrower ID", how="left")
if not dir_sum.empty:
    master = master.merge(dir_sum, on="Borrower ID", how="left")

print("Building presentation …")
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = NAVY

# Large brand band
band = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(4.5))
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()

# Teal accent strip
acc = slide.shapes.add_shape(1, Inches(0), Inches(4.5), SW, Pt(6))
acc.fill.solid()
acc.fill.fore_color.rgb = TEAL
acc.line.fill.background()

# Title
t = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11), Inches(1.6))
tf = t.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
r = p.add_run()
r.text = "FinnUp MSME Borrowers"
r.font.bold  = True
r.font.size  = Pt(40)
r.font.color.rgb = WHITE

# Sub-title
t2 = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(11), Inches(0.8))
tf2 = t2.text_frame
p2 = tf2.paragraphs[0]
r2 = p2.add_run()
r2.text = "Exploratory Data Analysis — Pre-Modelling Summary"
r2.font.size  = Pt(22)
r2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

# Date + context
t3 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11), Inches(0.5))
tf3 = t3.text_frame
p3 = tf3.paragraphs[0]
r3 = p3.add_run()
r3.text = "March 2026  |  6,483 Borrowers  |  11 Data Sheets"
r3.font.size  = Pt(14)
r3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

# Bottom note
t4 = slide.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11), Inches(1.5))
tf4 = t4.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
r4 = p4.add_run()
r4.text = ("Objective: Understand data quality, feature distributions and "
           "the label landscape across all source sheets before model development.\n"
           "Outcome labels are not yet available — ML kick-off planned once labels are received.")
r4.font.size  = Pt(12)
r4.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Agenda")
sections = [
    ("1", "Data Inventory",               "Row counts & null rates across all 11 sheets"),
    ("2", "Missing Value Analysis",       "Which columns are sparse vs. reliable"),
    ("3", "Borrower Profile",             "Company type, CIBIL, scale indicators"),
    ("4", "Proxy Labels",                 "Lead Status · FinnUp Status · Count Disbursed"),
    ("5", "Bank Statement Risk Signals",  "Bounce ratio, EOD balance, credit/debit ratio"),
    ("6", "Director Credit Profile",      "Director CIBIL distribution & comparison"),
    ("7", "Geographic Distribution",      "State-wise coverage & CIBIL variation"),
    ("8", "Financial KPIs",               "DSCR, current ratio, margins — 500 borrowers"),
    ("9", "Loan Products",                "Product mix & FinnUp conversion rate"),
    ("10","Feature Coverage for ML",      "What data is production-ready for modelling"),
    ("11","Key Findings & Next Steps",    "Action items before model development"),
]
col_x = [Inches(0.4), Inches(1.1), Inches(7.0)]
for i, (num, sec, desc) in enumerate(sections):
    y = Inches(1.35) + i * Inches(0.55)
    # number pill
    pill = slide.shapes.add_shape(1, col_x[0], y, Inches(0.55), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = TEAL if i % 2 == 0 else AMBER
    pill.line.fill.background()
    ptf = pill.text_frame
    pp  = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    pr  = pp.add_run()
    pr.text = num
    pr.font.bold  = True
    pr.font.size  = Pt(10)
    pr.font.color.rgb = WHITE
    # section name
    add_textbox(slide, col_x[1], y, Inches(5.5), Inches(0.4),
                sec, bold=True, size=11, color=NAVY)
    # description
    add_textbox(slide, col_x[2], y, Inches(6.0), Inches(0.4),
                desc, bold=False, size=10, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATA INVENTORY
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Data Inventory", "11 sheets loaded from FinnUp_Borrowers.xlsx")

sheet_info = []
for sh, df in data.items():
    if df.empty:
        sheet_info.append((sh.value, 0, 0, 100.0))
    else:
        null_pct = df.isnull().values.mean() * 100
        sheet_info.append((sh.value, len(df), len(df.columns), round(null_pct, 1)))

headers = ["Sheet", "Rows", "Cols", "Null %"]
col_widths = [Inches(3.2), Inches(1.3), Inches(1.0), Inches(1.2)]
col_starts = [Inches(0.4), Inches(3.7), Inches(5.1), Inches(6.2)]
row_h  = Inches(0.42)
top0   = Inches(1.35)

# header row
for j, (hdr, xp, wd) in enumerate(zip(headers, col_starts, col_widths)):
    hb = slide.shapes.add_shape(1, xp, top0, wd, row_h)
    hb.fill.solid()
    hb.fill.fore_color.rgb = NAVY
    hb.line.fill.background()
    htf = hb.text_frame
    htf.margin_left = Inches(0.07)
    htf.margin_top  = Inches(0.07)
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run()
    hr.text = hdr
    hr.font.bold  = True
    hr.font.size  = Pt(11)
    hr.font.color.rgb = WHITE

for i, (name, rows, cols, null_p) in enumerate(sheet_info):
    y   = top0 + row_h + i * row_h
    bg  = LIGHT if i % 2 == 0 else WHITE
    row_vals = [name, f"{rows:,}", str(cols), f"{null_p:.1f}%"]
    for j, (val, xp, wd) in enumerate(zip(row_vals, col_starts, col_widths)):
        cb = slide.shapes.add_shape(1, xp, y, wd, row_h)
        cb.fill.solid()
        cb.fill.fore_color.rgb = bg
        cb.line.fill.background()
        ctf = cb.text_frame
        ctf.margin_left = Inches(0.07)
        ctf.margin_top  = Inches(0.07)
        cp  = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.LEFT
        cr  = cp.add_run()
        cr.text = val
        cr.font.size  = Pt(10)
        # Colour null % red if > 50
        if j == 3 and null_p > 50:
            cr.font.color.rgb = SALMON
        else:
            cr.font.color.rgb = DARK

# Right-side stat boxes
stats = [
    ("6,483", "Total Borrowers"),
    ("11",    "Data Sheets"),
    ("30.3%", "Avg Null Rate\n(Profile sheet)"),
]
for k, (val, lbl) in enumerate(stats):
    add_stat_box(slide, Inches(7.8), Inches(1.55 + k * 1.6),
                 Inches(2.0), Inches(1.3), val, lbl,
                 bg=TEAL if k % 2 == 0 else NAVY)

add_bullet_box(slide, Inches(10.1), Inches(1.3), Inches(3.0), Inches(4.5),
               ["Financial Summary has 168 cols",
                "Financial KPIs: 36 KPI labels × multi-year",
                "Bank Statements: 3,996 rows / 2,489 borrowers",
                "Directors: 9,273 rows (1.4 per borrower avg)"],
               title="Highlights", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MISSING VALUE ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Missing Value Analysis — Borrower Profile",
                  "16 of 50 columns exceed the 50% missing threshold")

null_pct = (prof.isnull().sum() / len(prof) * 100).sort_values(ascending=False)
null_pct = null_pct[null_pct > 0]

fig, ax = plt.subplots(figsize=(11, 4.2), facecolor="white")
colors = ["#EF4444" if v > 50 else "#0D9488" for v in null_pct.values]
null_pct.plot(kind="bar", ax=ax, color=colors, edgecolor="white", width=0.8)
ax.axhline(50, color="red", linestyle="--", linewidth=1.2, label="50% threshold")
ax.set_title("Borrower Profile — % Missing per Column", fontsize=12, pad=8)
ax.set_ylabel("% Missing")
ax.set_xticklabels(ax.get_xticklabels(), rotation=40, ha="right", fontsize=8)
ax.legend(fontsize=9)
red_p  = mpatches.Patch(color="#EF4444", label=f"Above 50% ({(null_pct > 50).sum()} cols)")
teal_p = mpatches.Patch(color="#0D9488", label=f"Below 50% ({(null_pct <= 50).sum()} cols)")
ax.legend(handles=[red_p, teal_p], fontsize=9)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               ["CIBIL Score: ~99% missing",
                "Lead Status: ~99% missing",
                "CIBIL Score Range: 100% missing",
                "Company Website: 100% missing",
                "Platform / Source: ~99% missing",
                "Turnover: near-complete (100%)",
                "Networth: 92% coverage",
                "Company Name: 96% coverage",
                "State: 45% coverage",
                "Director CIBIL needed as proxy"],
               title="Coverage Notes", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BORROWER PROFILE
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Borrower Profile — Company Type & CIBIL",
                  "Proprietorships dominate; CIBIL data sparse (n=55) but skewed positive")

ct_col = "Company Type" if "Company Type" in prof.columns else "Company Type ID"
ct = prof[ct_col].value_counts(dropna=True).head(8)
cibil = prof["CIBIL Score"].dropna()

fig, axes = plt.subplots(1, 2, figsize=(11, 4.2), facecolor="white")
ct.plot(kind="bar", ax=axes[0], color="#1B3A6B", edgecolor="white", width=0.7)
axes[0].set_title("Company Type Distribution  (n=6,483)", fontsize=11)
axes[0].set_xticklabels(axes[0].get_xticklabels(), rotation=30, ha="right", fontsize=9)
axes[0].set_ylabel("Count")

axes[1].hist(cibil, bins=30, color="#0D9488", edgecolor="white", alpha=0.85)
axes[1].axvline(cibil.median(), color="red", linestyle="--",
                label=f"Median: {cibil.median():.0f}")
axes[1].axvline(700, color="orange", linestyle="--", label="Good threshold (700)")
axes[1].set_title(f"CIBIL Score Distribution  (n={len(cibil):,})", fontsize=11)
axes[1].set_xlabel("CIBIL Score")
axes[1].legend(fontsize=9)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Proprietorships: {ct.iloc[0]:,} ({ct.iloc[0]/len(prof)*100:.0f}%)",
                "Median CIBIL: 752  (positive skew)",
                "76% of CIBIL scores > 700 (where available)",
                "CIBIL coverage only 0.8% of borrowers",
                "Director CIBIL (86% coverage) must serve as proxy",
                "Company type useful as categorical feature"],
               title="Key Takeaways", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SCALE INDICATORS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "MSME Scale Indicators — Turnover & Networth",
                  "Both log-normally distributed; Turnover near-complete, Networth 92% covered")

fig, axes = plt.subplots(1, 2, figsize=(11, 4.2), facecolor="white")
for ax, col, color in [
    (axes[0], "Turnover",  "#1B3A6B"),
    (axes[1], "Networth",  "#F59E0B"),
]:
    if col not in prof.columns:
        ax.set_visible(False)
        continue
    vals = pd.to_numeric(prof[col], errors="coerce").dropna()
    vals = np.log1p(vals[vals > 0])
    ax.hist(vals, bins=40, color=color, edgecolor="white", alpha=0.85)
    ax.axvline(vals.median(), color="red", linestyle="--",
               label=f"Median log: {vals.median():.1f}")
    ax.set_title(f"{col}  (log scale, n={len(vals):,})", fontsize=11)
    ax.set_xlabel(f"log(1 + {col})")
    ax.legend(fontsize=9)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

t_vals = pd.to_numeric(prof.get("Turnover"), errors="coerce").dropna()
n_vals = pd.to_numeric(prof.get("Networth"), errors="coerce").dropna()
add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Turnover coverage: {len(t_vals):,} / 6,483 (100%)",
                f"Networth coverage: {len(n_vals):,} / 6,483 (92%)",
                "Both highly right-skewed → log transform required for ML",
                "Median log-Turnover ≈ 16.8 (≈ ₹20M range)",
                "Reliable features for model training",
                "Recommend clipping extreme outliers at 99th pct"],
               title="Key Takeaways", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PROXY LABELS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Proxy Labels — Target Variable Landscape",
                  "Outcome labels not yet received; best available proxies identified")

fig, axes = plt.subplots(1, 3, figsize=(11, 4.0), facecolor="white")

# Lead Status
if "Lead Status" in prof.columns:
    ls = prof["Lead Status"].fillna("No Label")
    ls.value_counts().plot(kind="bar", ax=axes[0],
                           color=["#1B3A6B","#0D9488","#F59E0B"][:len(ls.value_counts())],
                           edgecolor="white", width=0.6)
    axes[0].set_title(f"Lead Status  (n=6,483)", fontsize=10)
    axes[0].set_xticklabels(axes[0].get_xticklabels(), rotation=15, ha="right")

# Sanction vs Disbursed
s_col = pd.to_numeric(prof.get("Count Sanctions"), errors="coerce").dropna()
d_col = pd.to_numeric(prof.get("Count Disbursed"), errors="coerce").dropna()
axes[1].hist(s_col.clip(0, 5), bins=6, alpha=0.65, label=f"Sanctions (n={len(s_col):,})",
             color="#1B3A6B", edgecolor="white")
axes[1].hist(d_col.clip(0, 5), bins=6, alpha=0.65, label=f"Disbursed (n={len(d_col):,})",
             color="#F59E0B", edgecolor="white")
axes[1].set_title("Sanction vs Disbursement Counts", fontsize=10)
axes[1].legend(fontsize=8)
axes[1].set_xlim(-0.5, 5.5)

# FinnUp Status
if not loans.empty and "FinnUp Status" in loans.columns:
    fs = loans["FinnUp Status"].value_counts()
    fs.plot(kind="bar", ax=axes[2],
            color=["#1B3A6B","#EF4444"][:len(fs)],
            edgecolor="white", width=0.5)
    axes[2].set_title(f"FinnUp Status  (Loan Requests, n={len(loans):,})", fontsize=10)
    axes[2].set_xticklabels(axes[2].get_xticklabels(), rotation=0)

plt.suptitle("Proxy Label Overview", fontsize=11, y=1.01)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.4))

finnup_yes = (loans["FinnUp Status"] == "YES").sum() if (not loans.empty and "FinnUp Status" in loans.columns) else 0
finnup_no  = (loans["FinnUp Status"] == "NO").sum()  if (not loans.empty and "FinnUp Status" in loans.columns) else 0
add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               ["Lead Status: 99%+ have 'No Label' tag",
                "Only 3 labelled records (WARM / Hot Lead)",
                f"FinnUp Status YES: {finnup_yes:,} / {len(loans):,} requests ({finnup_yes/len(loans)*100:.1f}%)" if len(loans) else "FinnUp Status: N/A",
                "Count Disbursed > 0 for 399 borrowers (6%)",
                "FinnUp Status is the best available binary label",
                "True outcome labels (approved/rejected) awaited",
                "Target: labels expected within ~1 week"],
               title="Label Status", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BANK STATEMENT RISK SIGNALS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Bank Statement Risk Signals",
                  f"3,996 bank rows covering {bank['Borrower ID'].nunique():,} borrowers (38%)")

br  = bank["Bounce Ratio"].dropna().clip(0, 1)
eod = np.log1p(pd.to_numeric(bank["Avg EOD Balance"], errors="coerce").dropna().clip(lower=0))
eod = eod[np.isfinite(eod)]
cdr = bank["Credit Debit Ratio"].dropna().clip(0, 3)
ib  = pd.to_numeric(bank["Inward Cheque Bounces"], errors="coerce").clip(0, 20)
ob  = pd.to_numeric(bank["Outward Cheque Bounces"], errors="coerce").clip(0, 20)

fig, axes = plt.subplots(2, 2, figsize=(11, 5.5), facecolor="white")
axes[0,0].hist(br, bins=40, color="#EF4444", edgecolor="white", alpha=0.85)
axes[0,0].axvline(br.median(), color="black", linestyle="--",
                  label=f"Median: {br.median():.4f}")
axes[0,0].set_title("Bounce Ratio  (capped 1.0)", fontsize=10)
axes[0,0].legend(fontsize=8)

axes[0,1].hist(eod, bins=40, color="#10B981", edgecolor="white", alpha=0.85)
axes[0,1].set_title("Avg EOD Balance  (log scale)", fontsize=10)
axes[0,1].set_xlabel("log(1 + Balance)")

axes[1,0].scatter(ib, ob, alpha=0.25, s=8, color="#1B3A6B")
axes[1,0].set_xlabel("Inward Bounces (capped 20)")
axes[1,0].set_ylabel("Outward Bounces (capped 20)")
axes[1,0].set_title("Inward vs Outward Cheque Bounces", fontsize=10)

axes[1,1].hist(cdr, bins=40, color="#8B5CF6", edgecolor="white", alpha=0.85)
axes[1,1].axvline(1.0, color="red", linestyle="--", label="Credits = Debits")
axes[1,1].set_title("Credit / Debit Ratio", fontsize=10)
axes[1,1].legend(fontsize=8)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.2), Inches(9.0), Inches(5.6))

add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Unique borrowers with bank data: {bank['Borrower ID'].nunique():,}",
                "Median bounce ratio ≈ 0 (most accounts clean)",
                "Heavy tail in bounces — important risk outliers",
                "Credit/Debit ratio clusters tightly at 1.0",
                "EOD balance log-normally distributed",
                "Bank data covers only 38% of borrowers",
                "Key derived features: bounce_ratio, credit_debit_ratio, avg_eod_balance"],
               title="Key Takeaways", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DIRECTOR CREDIT PROFILE
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Director Credit Profile",
                  f"9,273 director records across 5,600 unique borrowers (86% coverage)")

dcibil = dirs["CIBIL Score"].dropna()
fig, axes = plt.subplots(1, 2, figsize=(11, 4.2), facecolor="white")
axes[0].hist(dcibil, bins=35, color="#1B3A6B", edgecolor="white", alpha=0.85)
axes[0].axvline(dcibil.median(), color="red", linestyle="--",
                label=f"Median: {dcibil.median():.0f}")
axes[0].axvline(700, color="orange", linestyle="--", label="700 threshold")
axes[0].set_title(f"Director CIBIL Distribution  (n={len(dcibil):,})", fontsize=11)
axes[0].legend(fontsize=9)

m = (prof[["Borrower ID","CIBIL Score"]]
     .merge(dir_sum[["Borrower ID","min_director_cibil"]], on="Borrower ID", how="inner")
     .dropna())
if len(m) > 0:
    axes[1].scatter(m["CIBIL Score"], m["min_director_cibil"],
                    alpha=0.4, s=12, color="#0D9488")
    axes[1].plot([300,900],[300,900],"r--", linewidth=1, label="Equal line")
    axes[1].set_xlabel("Borrower CIBIL")
    axes[1].set_ylabel("Min Director CIBIL")
    axes[1].set_title(f"Borrower vs Min Director CIBIL  (n={len(m):,})", fontsize=11)
    axes[1].legend(fontsize=9)
else:
    axes[1].text(0.5, 0.5, "Insufficient overlap data", ha="center", va="center",
                 transform=axes[1].transAxes, fontsize=12, color="gray")
    axes[1].set_title("Borrower vs Min Director CIBIL", fontsize=11)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Director CIBIL n = {len(dcibil):,} records",
                f"Median Director CIBIL: {dcibil.median():.0f}",
                f"{(dcibil >= 700).sum() / len(dcibil) * 100:.0f}% directors score ≥ 700",
                "Avg 1.4 directors per borrower",
                "min_director_cibil is a strong risk feature",
                "86% borrower coverage — best credit proxy available"],
               title="Key Takeaways", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GEOGRAPHIC DISTRIBUTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Geographic Distribution",
                  "28 states covered; Maharashtra & Tamil Nadu dominate; 45% of borrowers have state data")

fig, axes = plt.subplots(1, 2, figsize=(11, 4.8), facecolor="white")
sc = prof["State"].value_counts(dropna=True).head(15)
sc.plot(kind="bar", ax=axes[0], color="#1B3A6B", edgecolor="white", width=0.7)
axes[0].set_title("Top 15 States by Borrower Count", fontsize=11)
axes[0].set_xticklabels(axes[0].get_xticklabels(), rotation=40, ha="right", fontsize=8)
axes[0].set_ylabel("Count")

top8 = prof["State"].value_counts().head(8).index
sub  = prof[prof["State"].isin(top8)][["State","CIBIL Score"]].dropna()
if not sub.empty:
    order = sub.groupby("State")["CIBIL Score"].median().sort_values(ascending=False).index.tolist()
    sns.boxplot(data=sub, x="State", y="CIBIL Score", order=order,
                ax=axes[1], color="#0D9488")
    axes[1].set_title(f"CIBIL by State — Top 8  (n={len(sub):,})", fontsize=11)
    axes[1].tick_params(axis="x", rotation=30, labelsize=8)
    axes[1].set_xlabel("")
else:
    axes[1].text(0.5, 0.5, "Insufficient CIBIL overlap", ha="center", va="center",
                 transform=axes[1].transAxes, fontsize=12, color="gray")
    axes[1].set_title("CIBIL by State", fontsize=11)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

top3 = prof["State"].value_counts().head(3)
add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"#1 Maharashtra: {top3.iloc[0]:,} borrowers",
                f"#2 Tamil Nadu: {top3.iloc[1]:,} borrowers",
                f"#3 {top3.index[2]}: {top3.iloc[2]:,} borrowers",
                "State coverage: 45% (2,926 / 6,483)",
                "State useful as region-level feature",
                "Telangana shows highest median CIBIL",
                "State + pincode can proxy economic zone"],
               title="Key Takeaways", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FINANCIAL KPIs
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Financial KPIs — Latest Year Available",
                  f"500 borrowers have KPI data (7.7%); extreme outliers require treatment")

available = [k for k in key_kpis if k in kpi_pivot.columns]
fig, axes = plt.subplots(2, 3, figsize=(11, 5.5), facecolor="white")
axes = axes.flatten()
for ax, kpi in zip(axes, available):
    vals = kpi_pivot[kpi].dropna()
    p1, p99 = vals.quantile(0.01), vals.quantile(0.99)
    ax.hist(vals.clip(p1, p99), bins=25, color="#1B3A6B", edgecolor="white", alpha=0.85)
    ax.axvline(vals.median(), color="#EF4444", linestyle="--",
               label=f"Median: {vals.median():.2f}")
    ax.set_title(f"{kpi}", fontsize=8.5)
    ax.legend(fontsize=7)
    ax.tick_params(labelsize=7)
for ax in axes[len(available):]:
    ax.set_visible(False)
plt.suptitle("KPI Distributions — 1st–99th pct clipped", fontsize=10, y=1.01)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.2), Inches(9.0), Inches(5.6))

kpi_med = {k: kpi_pivot[k].median() for k in available if k in kpi_pivot.columns}
add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Borrowers with KPI data: {len(kpi_pivot):,} / 6,483",
                f"Median DSCR: {kpi_med.get('DSCR (Avg/Min)', 0):.2f}  (>1.0 = healthy)",
                f"Median Current Ratio: {kpi_med.get('Current Ratio', 0):.2f}",
                f"Median Debt/Equity: {kpi_med.get('Debt Equity Ratio', 0):.2f}",
                f"Median Net Profit Margin: {kpi_med.get('Net Profit Margin (%)', 0):.2f}%",
                "High variance — extreme outliers present",
                "Clip at 1–99th percentile before training"],
               title="KPI Summary", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — LOAN PRODUCTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Loan Product Mix & FinnUp Conversion",
                  "Unsecured Business Loan dominates; overall FinnUp Status YES rate < 5%")

fig, axes = plt.subplots(1, 2, figsize=(11, 4.5), facecolor="white")
if not loans.empty and "Product Name" in loans.columns:
    top12 = loans["Product Name"].value_counts().head(12)
    top12.plot(kind="barh", ax=axes[0], color="#1B3A6B", edgecolor="white")
    axes[0].set_title("Top 12 Loan Products Applied", fontsize=11)
    axes[0].invert_yaxis()
    axes[0].tick_params(labelsize=8)

if not loans.empty and {"Product Name","FinnUp Status"}.issubset(loans.columns):
    top8 = loans["Product Name"].value_counts().head(8).index
    ct = pd.crosstab(loans[loans["Product Name"].isin(top8)]["Product Name"],
                     loans["FinnUp Status"], normalize="index") * 100
    ct.plot(kind="bar", stacked=True, ax=axes[1],
            color=["#1B3A6B","#EF4444"], edgecolor="white")
    axes[1].set_title("FinnUp Status by Product (Top 8)", fontsize=11)
    axes[1].set_ylabel("% of Applications")
    axes[1].set_xticklabels(axes[1].get_xticklabels(), rotation=35, ha="right", fontsize=8)
    axes[1].legend(title="FinnUp Status", fontsize=9)
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(9.0), Inches(5.5))

finnup_rate = finnup_yes / len(loans) * 100 if len(loans) else 0
add_bullet_box(slide, Inches(9.5), Inches(1.35), Inches(3.6), Inches(5.0),
               [f"Total loan requests: {len(loans):,}",
                f"FinnUp YES: {finnup_yes:,} ({finnup_rate:.1f}%)",
                "Top product: Unsecured Business Loan",
                "Personal Loan #2, Cash Credit #3",
                "Unsecured Business Loan has highest YES rate",
                "Product type is a useful categorical feature",
                "Class imbalance needs SMOTE / class weights"],
               title="Product Insights", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FEATURE COVERAGE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Feature Coverage for ML",
                  "Data quality assessment — what is ready for training today")

def col_cov(df, col):
    if col not in df.columns: return 0
    return int(pd.to_numeric(df[col], errors="coerce").notna().sum())

feature_groups = [
    ("Turnover",                   col_cov(prof, "Turnover"),        6483),
    ("Networth",                   col_cov(prof, "Networth"),         6483),
    ("Director CIBIL (min)",       int(dir_sum["Borrower ID"].nunique()) if not dir_sum.empty else 0, 6483),
    ("State / Location",           int(prof["State"].notna().sum()) if "State" in prof.columns else 0, 6483),
    ("Bank Account Data",          int(bank["Borrower ID"].nunique()) if not bank.empty else 0, 6483),
    ("Financial KPIs",             int(kpi_pivot["Borrower ID"].nunique()) if not kpi_pivot.empty else 0, 6483),
    ("CIBIL Score (borrower)",     col_cov(prof, "CIBIL Score"),     6483),
    ("Count Disbursed > 0 (label)",int((pd.to_numeric(prof.get("Count Disbursed"), errors="coerce") > 0).sum()), 6483),
    ("FinnUp Status (label)",      finnup_yes,                        len(loans)),
]

fig, ax = plt.subplots(figsize=(10, 4.2), facecolor="white")
names   = [f[0] for f in feature_groups]
pcts    = [f[1] / f[2] * 100 for f in feature_groups]
counts  = [f[1] for f in feature_groups]
colors  = ["#10B981" if p >= 80 else "#F59E0B" if p >= 40 else "#EF4444" for p in pcts]
bars    = ax.barh(names, pcts, color=colors, edgecolor="white", height=0.6)
ax.axvline(80, color="#10B981", linestyle="--", linewidth=1.2, label="80% (high coverage)")
ax.axvline(40, color="#F59E0B", linestyle="--", linewidth=1.2, label="40% (medium coverage)")
for bar, cnt, pct in zip(bars, counts, pcts):
    ax.text(pct + 1, bar.get_y() + bar.get_height() / 2,
            f"{pct:.1f}%  ({cnt:,})", va="center", fontsize=9)
ax.set_xlim(0, 115)
ax.set_xlabel("Coverage %")
ax.set_title("Feature Data Coverage (% of 6,483 Borrowers or Loan Requests)", fontsize=11)
ax.legend(fontsize=9)
ax.invert_yaxis()
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.2), Inches(9.3), Inches(5.6))

add_bullet_box(slide, Inches(9.7), Inches(1.35), Inches(3.4), Inches(5.0),
               ["Green (≥80%): Ready for training",
                "Amber (40-80%): Use with imputation",
                "Red (<40%): Treat as sparse signal",
                "─────────────────────",
                "Turnover & Networth: production-ready",
                "Director CIBIL: key proxy for credit risk",
                "Bank data: sparse but high-signal",
                "Label (FinnUp YES): highly imbalanced"],
               title="Coverage Legend", size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — KEY FINDINGS & NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Key Findings & Next Steps",
                  "Summary of EDA — actions before model development")

findings = [
    ("Data is rich but sparse in places",
     "Turnover (100%) and Networth (92%) are reliable. CIBIL almost absent; Director CIBIL @ 86% is the credit proxy."),
    ("Labels are the bottleneck",
     "Only ~5% of loan requests have FinnUp Status YES. True approved/rejected/defaulted labels are awaited."),
    ("Bank data is high-signal, low-coverage",
     "Bounce ratio, credit/debit ratio and EOD balance are excellent risk features — but cover only 38% of borrowers."),
    ("Severe class imbalance expected",
     "FinnUp YES rate ~5%. ML model needs SMOTE, class weighting, or threshold tuning."),
    ("Company type & geography add context",
     "Proprietorships dominate. Maharashtra + Tamil Nadu cover ~27% of volume. State is a useful regional proxy."),
    ("Financial KPIs need careful treatment",
     "500-borrower coverage and extreme outliers. Clip to 1–99th percentile; build separate KPI-enriched model path."),
]

y_start = Inches(1.3)
for i, (title, body) in enumerate(findings):
    row_y = y_start + i * Inches(0.98)
    icon = slide.shapes.add_shape(1, Inches(0.3), row_y, Inches(0.35), Inches(0.35))
    icon.fill.solid()
    icon.fill.fore_color.rgb = TEAL if i % 2 == 0 else AMBER
    icon.line.fill.background()
    itf = icon.text_frame
    itf.margin_top = Inches(0.02)
    ip  = itf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    ir  = ip.add_run()
    ir.text = str(i + 1)
    ir.font.bold  = True
    ir.font.size  = Pt(10)
    ir.font.color.rgb = WHITE

    add_textbox(slide, Inches(0.78), row_y, Inches(6.1), Inches(0.45),
                title, bold=True, size=11, color=NAVY)
    add_textbox(slide, Inches(0.78), row_y + Inches(0.38), Inches(6.1), Inches(0.5),
                body, bold=False, size=9.5, color=GRAY)

# Next steps box
ns_box = slide.shapes.add_shape(1, Inches(7.1), Inches(1.3), Inches(5.9), Inches(5.8))
ns_box.fill.solid()
ns_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
ns_box.line.color.rgb = NAVY
ns_box.line.width = Pt(1.5)

add_textbox(slide, Inches(7.25), Inches(1.4), Inches(5.7), Inches(0.45),
            "NEXT STEPS", bold=True, size=13, color=NAVY)

next_steps = [
    ("🏷️  This week",     "Receive outcome labels from FinnUp team"),
    ("🔧  Week 2",        "Run 02_feature_engineering.ipynb — encode, impute, derive features"),
    ("🤖  Week 2–3",      "Train baseline models (LightGBM + Logistic Regression)"),
    ("📊  Week 3",        "Cross-validate; tune for precision/recall trade-off"),
    ("🔍  Week 3–4",      "SHAP explainability + lender policy integration"),
    ("🚀  Week 4+",       "Scoring pipeline per lender product"),
]
for k, (phase, task) in enumerate(next_steps):
    y_ns = Inches(1.95) + k * Inches(0.82)
    pb = slide.shapes.add_shape(1, Inches(7.25), y_ns, Inches(1.2), Inches(0.32))
    pb.fill.solid()
    pb.fill.fore_color.rgb = NAVY if k == 0 else TEAL if k % 2 else AMBER
    pb.line.fill.background()
    ptf2 = pb.text_frame
    ptf2.margin_left = Inches(0.05)
    ptf2.margin_top  = Inches(0.04)
    pp2  = ptf2.paragraphs[0]
    pp2.alignment = PP_ALIGN.CENTER
    pr2  = pp2.add_run()
    pr2.text = phase
    pr2.font.bold  = True
    pr2.font.size  = Pt(8)
    pr2.font.color.rgb = WHITE
    add_textbox(slide, Inches(8.55), y_ns, Inches(4.3), Inches(0.32),
                task, bold=False, size=9.5, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU / FOOTER
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = NAVY

acc = slide.shapes.add_shape(1, Inches(0), Inches(3.2), SW, Pt(4))
acc.fill.solid()
acc.fill.fore_color.rgb = TEAL
acc.line.fill.background()

t = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1.2))
tf = t.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "FinnUp MSME — EDA Complete"
r.font.bold  = True
r.font.size  = Pt(34)
r.font.color.rgb = WHITE

t2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.75), Inches(10), Inches(0.6))
tf2 = t2.text_frame
p2  = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2  = p2.add_run()
r2.text = "Awaiting outcome labels · Next: Feature Engineering & Model Training"
r2.font.size  = Pt(15)
r2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

t3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(2.5))
tf3 = t3.text_frame
tf3.word_wrap = True
p3  = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3  = p3.add_run()
r3.text = ("Data  ·  6,483 Borrowers  ·  11 Sheets  ·  March 2026\n"
           "Repo: finnup-prediction / notebooks / 01_eda.ipynb")
r3.font.size  = Pt(12)
r3.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

# ── SAVE ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\n✅  Saved: {OUTPUT_PATH.resolve()}")
print(f"   Slides: {len(prs.slides)}")
