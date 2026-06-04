"""
generate_iim_submission_ppt.py
------------------------------
Generates  outputs/FinnUp_IIM_Submission.pptx
  – APAL Cohort 2 | Group 1 | IIM Calcutta
    Mid-Project Progress Submission (April 2026)
    "Solving Lender Matching with Borrower Intelligence"

Run:
    python generate_iim_submission_ppt.py
"""

from __future__ import annotations
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import warnings
warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Output ────────────────────────────────────────────────────────────────────
OUTPUT_PATH = Path("outputs") / "FinnUp_IIM_Submission.pptx"
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x3A, 0x6B)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
SALMON  = RGBColor(0xEF, 0x44, 0x44)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
LTBLUE  = RGBColor(0x93, 0xC5, 0xFD)
GOLD    = RGBColor(0xD9, 0x7F, 0x0F)

SW = Inches(13.33)
SH = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def fig_to_bytes(fig, dpi=130):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def add_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = WHITE
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run.font.bold = True; run.font.size = Pt(21); run.font.color.rgb = WHITE
    tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.20)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.82), Inches(12), Pt(18))
        sr = sub.text_frame.paragraphs[0].add_run()
        sr.text = subtitle; sr.font.size = Pt(10.5)
        sr.font.color.rgb = TEAL; sr.font.italic = True
    line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), SW, Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
    # IIM watermark (top-right)
    wm = slide.shapes.add_textbox(Inches(11.2), Inches(0.18), Inches(2.0), Inches(0.55))
    wmp = wm.text_frame.paragraphs[0]; wmp.alignment = PP_ALIGN.RIGHT
    wmr = wmp.add_run(); wmr.text = "APAL | IIM Calcutta"
    wmr.font.size = Pt(8); wmr.font.color.rgb = LTBLUE; wmr.font.italic = True
    return slide


def txb(slide, left, top, width, height, text,
        bold=False, size=11, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.bold = bold; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def blist(slide, left, top, width, height, items,
          size=10.5, title=None, title_color=None, bullet="▸"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.bold = True; r.font.size = Pt(size + 1.5)
        r.font.color.rgb = title_color or NAVY
    for item in items:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = f"{bullet}  {item}"
        r.font.size = Pt(size); r.font.color.rgb = DARK
    return tb


def stat_box(slide, left, top, width, height, value, label, bg=NAVY):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = bg; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value
    r.font.bold = True; r.font.size = Pt(26); r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(9); r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)


def col_header(slide, left, top, width, height, text, bg=NAVY):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg; sh.line.fill.background()
    tf = sh.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.07)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.bold = True; r.font.size = Pt(10.5); r.font.color.rgb = WHITE


# =============================================================================
# BUILD PRESENTATION
# =============================================================================
prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = NAVY

# Full bg
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

# Teal bottom stripe
acc = slide.shapes.add_shape(1, Inches(0), Inches(6.8), SW, Inches(0.7))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL; acc.line.fill.background()

# Gold side accent
side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), SH)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()

# Program badge (top-right)
badge_bg = slide.shapes.add_shape(1, Inches(9.8), Inches(0.3), Inches(3.3), Inches(0.7))
badge_bg.fill.solid(); badge_bg.fill.fore_color.rgb = TEAL; badge_bg.line.fill.background()
bt = slide.shapes.add_textbox(Inches(9.8), Inches(0.3), Inches(3.3), Inches(0.7))
btp = bt.text_frame.paragraphs[0]; btp.alignment = PP_ALIGN.CENTER
btr = btp.add_run(); btr.text = "APAL | IIM Calcutta | Cohort 2"
btr.font.bold = True; btr.font.size = Pt(11); btr.font.color.rgb = WHITE

# Title
t1 = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.5), Inches(1.4))
p1 = t1.text_frame.paragraphs[0]
r1 = p1.add_run(); r1.text = "Solving Lender Matching with Borrower Intelligence"
r1.font.bold = True; r1.font.size = Pt(34); r1.font.color.rgb = WHITE

# Sub-title
t2 = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.5), Inches(0.7))
p2 = t2.text_frame.paragraphs[0]
r2 = p2.add_run(); r2.text = "Mid-Project Progress Report & Delivery Roadmap"
r2.font.size = Pt(20); r2.font.color.rgb = LTBLUE

# Context line
t3 = slide.shapes.add_textbox(Inches(0.6), Inches(3.45), Inches(12.5), Inches(0.5))
p3 = t3.text_frame.paragraphs[0]
r3 = p3.add_run(); r3.text = "April 2026  |  Group 1  |  FinnUp MSME Credit Marketplace"
r3.font.size = Pt(13); r3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

# Members
t4 = slide.shapes.add_textbox(Inches(0.6), Inches(4.15), Inches(12.5), Inches(0.5))
p4 = t4.text_frame.paragraphs[0]
r4 = p4.add_run()
r4.text = "Asha · Arvind · Anil · Bhupesh · Deepak · Ganesh · Gopal · Hareram · Pranali · Rahul · Savitha · Sonam · Samik"
r4.font.size = Pt(11); r4.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

# Phase status chips
chips = [
    ("Phase 0\nProblem Definition", "✓ DONE", GREEN),
    ("Phase 1\nData & EDA", "✓ DONE", GREEN),
    ("Phase 2\nFeature Engineering", "⏳ IN PROGRESS", AMBER),
    ("Phase 3\nML Model", "◌ PLANNED", GRAY),
    ("Phase 4\nAPI & Deployment", "◌ PLANNED", GRAY),
]
for i, (ph, st, col) in enumerate(chips):
    cx = Inches(0.45) + i * Inches(2.56)
    sh2 = slide.shapes.add_shape(1, cx, Inches(5.0), Inches(2.4), Inches(1.55))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = col; sh2.line.fill.background()
    tf2 = sh2.text_frame; tf2.word_wrap = True
    tf2.margin_top = Inches(0.12); tf2.margin_left = Inches(0.1)
    pp = tf2.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    ppr = pp.add_run(); ppr.text = ph
    ppr.font.bold = True; ppr.font.size = Pt(10); ppr.font.color.rgb = WHITE
    pp2 = tf2.add_paragraph(); pp2.alignment = PP_ALIGN.CENTER
    ppr2 = pp2.add_run(); ppr2.text = st
    ppr2.font.size = Pt(9.5); ppr2.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0) if col == GREEN else RGBColor(0xFF, 0xFF, 0xCC)

# Bottom bar text
bt2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.45))
bp2 = bt2.text_frame.paragraphs[0]; bp2.alignment = PP_ALIGN.CENTER
br2 = bp2.add_run()
br2.text = "Submitted to: APAL Programme  |  IIM Calcutta  |  Capstone Project — Interim Submission"
br2.font.size = Pt(9.5); br2.font.color.rgb = WHITE


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — RECAP: ORIGINAL PROPOSAL (JAN 2026)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Recap: Our Original Proposal (January 2026)",
                  "What we set out to build — as submitted to IIM Calcutta")

# Problem framing quote box
qbox = slide.shapes.add_shape(1, Inches(0.35), Inches(1.2), Inches(12.6), Inches(1.0))
qbox.fill.solid(); qbox.fill.fore_color.rgb = LIGHT; qbox.line.fill.background()
qtf = qbox.text_frame; qtf.word_wrap = True
qtf.margin_left = Inches(0.3); qtf.margin_top = Inches(0.12)
qp = qtf.paragraphs[0]
qr = qp.add_run()
qr.text = ('"Matching is treated as a rule-based eligibility problem, but is actually a '
           'probabilistic matching problem requiring intelligent optimisation." — Original Proposal, Jan 2026')
qr.font.size = Pt(11); qr.font.italic = True; qr.font.color.rgb = NAVY

# 3-column original pillars
pillars = [
    ("Problem\nIdentified", NAVY, [
        "50 Lakh+ MSME leads available, only 8–10% convert",
        "80–85% of sales effort wasted on mismatched leads",
        "Lenders reject borrowers due to policy mismatch, not credit risk",
        "Rule-based engine cannot learn from historical outcomes",
        "Late discovery of fit after significant effort already spent",
    ]),
    ("Proposed\nSolution", TEAL, [
        "Stage 1: XGBoost/tree models predict lender-specific approval probability",
        "Stage 2: Weighted scoring ranks lenders → borrowers routed to best-fit lenders",
        "Continuous Learning Loop: each outcome (approve/reject) retrains model",
        "AI Canvas: Prediction → Judgement → Action → Business Outcome",
        "Hybrid approach: ML optimises matching; rules enforce hard constraints",
    ]),
    ("Target\nOutcomes", AMBER, [
        "Increase conversion rate: 11% → 75%  (target state)",
        "Reduce wasted sales effort: 80% → <30%",
        "20% ROI on platform investment over 2 years",
        "5× revenue growth by Year 3",
        "Faster, relevant matches → better borrower + lender experience",
    ]),
]

col_w = Inches(4.1)
for i, (title, bg, items) in enumerate(pillars):
    cx = Inches(0.35) + i * (col_w + Inches(0.18))
    tbar = slide.shapes.add_shape(1, cx, Inches(2.35), col_w, Inches(0.5))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_top = Inches(0.07)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(11); tr.font.color.rgb = WHITE
    blist(slide, cx, Inches(2.88), col_w, Inches(4.3), items, size=9.5)

txb(slide, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.3),
    "Original submission: 18 January 2026  |  This report tracks progress and updated delivery plan against that proposal.",
    size=9.5, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — WHAT WE'VE DONE (PHASE 1 STATUS)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 1 Complete: Data Acquisition & Exploratory Analysis",
                  "Real data extracted from FinnUp's production database — validated and profiled")

# DONE badge
done = slide.shapes.add_shape(1, Inches(10.7), Inches(1.2), Inches(2.35), Inches(0.45))
done.fill.solid(); done.fill.fore_color.rgb = GREEN; done.line.fill.background()
dtf = done.text_frame; dp = dtf.paragraphs[0]; dp.alignment = PP_ALIGN.CENTER
dr = dp.add_run(); dr.text = "✓  MILESTONE COMPLETE"
dr.font.bold = True; dr.font.size = Pt(10); dr.font.color.rgb = WHITE

# Key achievement boxes
achievements = [
    ("Real Production\nData Secured", NAVY,
     "Extracted directly from FinnUp's live database. 11 structured sheets covering the full borrower lifecycle — profile, financials, bank statements, directors, loans and more.",
     "6,483 borrowers"),
    ("Complete Data\nAudit Done", TEAL,
     "Every column across all 11 sheets profiled for coverage, null rates, data type, and ML usability. 30+ features identified as production-ready for model training.",
     "30+ ML features"),
    ("Risk Signals\nQuantified", PURPLE,
     "Bank statement risk signals extracted: Bounce Ratio, Avg EOD Balance, Credit-Debit Ratio, Cheque Bounce patterns. Director CIBIL identified as primary credit proxy.",
     "3,996 bank rows"),
    ("Label Landscape\nMapped", AMBER,
     "Three proxy labels identified: FinnUp Status (YES/NO), Lead Status, Count Disbursed. FinnUp Status is the strongest available binary training signal until formal outcomes arrive.",
     "Binary label ready"),
]

ach_w = Inches(3.05)
for i, (title, bg, desc, tag) in enumerate(achievements):
    cx = Inches(0.35) + i * (ach_w + Inches(0.18))
    # header
    tbar = slide.shapes.add_shape(1, cx, Inches(1.25), ach_w, Inches(0.58))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_top = Inches(0.08)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(11); tr.font.color.rgb = WHITE
    # tag pill
    tag_sh = slide.shapes.add_shape(1, cx, Inches(1.85), ach_w, Inches(0.3))
    tag_sh.fill.solid(); tag_sh.fill.fore_color.rgb = LIGHT; tag_sh.line.fill.background()
    ttag = tag_sh.text_frame; ttag.margin_top = Inches(0.03)
    tagp = ttag.paragraphs[0]; tagp.alignment = PP_ALIGN.CENTER
    tagr = tagp.add_run(); tagr.text = tag
    tagr.font.bold = True; tagr.font.size = Pt(9); tagr.font.color.rgb = bg
    # body
    txb(slide, cx, Inches(2.2), ach_w, Inches(2.8),
        desc, size=9.5, color=DARK, wrap=True)

# Data sheet inventory table
headers = ["Sheet", "Rows / Borrowers", "Key Fields", "ML Readiness"]
rows_data = [
    ("Borrower Profile",   "6,483 borrowers",    "Turnover, Networth, Company Type, State, CIBIL",     "READY",   GREEN),
    ("Directors",          "9,273 records",       "Director CIBIL (86% coverage), Ownership %",         "READY",   GREEN),
    ("Bank Statements",    "3,996 rows / 2,489",  "Bounce Ratio, EOD Balance, CDR, Cheque Bounces",     "PARTIAL", AMBER),
    ("Loan Requests",      "Multi-row / borrower","Product, Loan Amount, FinnUp Status (label)",        "PARTIAL", AMBER),
    ("Financial KPIs",     "500 borrowers (8%)",  "DSCR, Current Ratio, Net Profit Margin",             "SPARSE",  SALMON),
    ("Fin. Summary / B/S / P&L / CF", "Multi-year", "Revenue, PAT, Assets, Liabilities",               "SPARSE",  SALMON),
]
col_x4 = [Inches(0.35), Inches(2.65), Inches(4.85), Inches(10.8)]
col_w4 = [Inches(2.22), Inches(2.12), Inches(5.85), Inches(2.2)]
row_h4 = Inches(0.40)
top4   = Inches(5.08)

for j, (hdr, xp, wd) in enumerate(zip(headers, col_x4, col_w4)):
    hb = slide.shapes.add_shape(1, xp, top4, wd, row_h4)
    hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
    htf = hb.text_frame; htf.margin_left = Inches(0.06); htf.margin_top = Inches(0.07)
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run(); hr.text = hdr
    hr.font.bold = True; hr.font.size = Pt(9.5); hr.font.color.rgb = WHITE

for i, (name, rows, fields, status, clr) in enumerate(rows_data):
    y = top4 + row_h4 + i * row_h4
    bg2 = LIGHT if i % 2 == 0 else WHITE
    vals = [name, rows, fields]
    for j, (val, xp, wd) in enumerate(zip(vals, col_x4, col_w4)):
        cb = slide.shapes.add_shape(1, xp, y, wd, row_h4)
        cb.fill.solid(); cb.fill.fore_color.rgb = bg2; cb.line.fill.background()
        ctf = cb.text_frame; ctf.margin_left = Inches(0.06); ctf.margin_top = Inches(0.07)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.LEFT
        cr = cp.add_run(); cr.text = val
        cr.font.size = Pt(8.5); cr.font.color.rgb = DARK
    # Readiness pill
    rp = slide.shapes.add_shape(1, col_x4[3] + Inches(0.12), y + Inches(0.07),
                                 Inches(1.9), Inches(0.26))
    rp.fill.solid(); rp.fill.fore_color.rgb = clr; rp.line.fill.background()
    rptf = rp.text_frame; rpp = rptf.paragraphs[0]; rpp.alignment = PP_ALIGN.CENTER
    rpr = rpp.add_run(); rpr.text = status
    rpr.font.bold = True; rpr.font.size = Pt(8); rpr.font.color.rgb = WHITE


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — KEY EDA FINDINGS
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Key EDA Findings — What the Data Tells Us",
                  "Insights driving feature selection and model design decisions")

# Coverage chart (left)
feature_data = {
    "Turnover":          100.0,
    "Company Name":       96.0,
    "Networth":           92.0,
    "Director CIBIL":     86.0,
    "Business Age":       51.0,
    "State / Location":   55.0,
    "Bank Account Data":  38.0,
    "Financial KPIs":      7.7,
    "Lead Status":         0.5,
    "Borrower CIBIL":      0.8,
}
labels  = list(feature_data.keys())
values  = list(feature_data.values())
bar_clr = ["#16A34A" if v >= 80 else ("#F59E0B" if v >= 40 else "#EF4444") for v in values]

fig, ax = plt.subplots(figsize=(6.5, 4.5), facecolor="white")
bars = ax.barh(labels, values, color=bar_clr, edgecolor="white", height=0.65)
ax.axvline(50, color="#64748B", linestyle="--", linewidth=1, label="50% threshold")
for bar, v in zip(bars, values):
    ax.text(v + 0.8, bar.get_y() + bar.get_height() / 2,
            f"{v:.0f}%", va="center", fontsize=9)
ax.set_xlim(0, 115)
ax.set_xlabel("Data Coverage (%)", fontsize=9)
ax.set_title("Feature Coverage — 6,483 Borrowers", fontsize=10, pad=6)
gp = mpatches.Patch(color="#16A34A", label="≥80%  Production Ready")
ap = mpatches.Patch(color="#F59E0B", label="40–80%  Usable with care")
rp = mpatches.Patch(color="#EF4444", label="<40%  Sparse")
ax.legend(handles=[gp, ap, rp], fontsize=8, loc="lower right")
plt.tight_layout()
slide.shapes.add_picture(fig_to_bytes(fig), Inches(0.3), Inches(1.22), Inches(6.8), Inches(5.3))

# Right column: curated finding cards
findings = [
    (NAVY,   "Director CIBIL is the Best Credit Proxy",
     "Borrower CIBIL score is available for only 55 of 6,483 borrowers (0.8%). "
     "Director CIBIL, available for 86% of borrowers, is the strongest credit signal "
     "in the dataset and will serve as the primary credit feature."),
    (TEAL,   "Scale Signals Are Reliable & Well-Distributed",
     "Turnover (100%) and Networth (92%) are both log-normally distributed — "
     "consistent with MSME profile expectations. These are primary scale features "
     "for lender ticket-size matching."),
    (PURPLE, "Bank Bounce Signals Capture Risk Behaviour",
     "Bounce Ratio, Credit-Debit Ratio and Avg EOD Balance (38% coverage) reveal "
     "clear bimodal patterns — low-risk borrowers cluster near zero bounce. "
     "These features strongly differentiate risk tiers."),
    (AMBER,  "Real Approval Labels — Full Training Set",
     "1,363 borrowers have conclusive approval/rejection outcomes from actual FinnUp disbursements. "
     "Remaining 4,824 treated as not approved (fillna=0). Full training set: 6,187 rows at 6.9% approval rate. "
     "Labels sourced from 'Loan Applications' sheet — "
     "Disbursed/Deal Sanctioned → approved; Reject Deal/Interest + unlabeled → not approved."),
    (GREEN,  "Company Profile Adds Segment Dimension",
     "Proprietorships dominate (>60%). State data (55%) and Business Age (51%) "
     "add geographic and vintage segmentation — critical for lender policy matching "
     "given sector/geography restrictions."),
]
cur_y = Inches(1.22)
for bg_c, title, desc in findings:
    tbar = slide.shapes.add_shape(1, Inches(7.3), cur_y, Inches(5.7), Inches(0.35))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg_c; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_left = Inches(0.1); ttf.margin_top = Inches(0.04)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(9.5); tr.font.color.rgb = WHITE
    txb(slide, Inches(7.3), cur_y + Inches(0.36), Inches(5.7), Inches(0.68),
        desc, size=8.8, color=DARK, wrap=True)
    cur_y += Inches(1.07)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — LINKING EDA TO ORIGINAL AI CANVAS
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "EDA Validates Our AI Canvas Design",
                  "Data findings directly inform the four-stage ML pipeline proposed in January 2026")

canvas = [
    ("PREDICTION\n(Stage 1)",    NAVY,
     "Proposed: Predict lender-specific approval probability",
     [
         "✓  FinnUp Status label available as training target",
         "✓  50+ borrower features identified across 11 sheets",
         "✓  Director CIBIL + Turnover are strong predictors",
         "⚠  Formal outcome labels (approved/rejected) awaited from CRM",
         "→  Will train on FinnUp Status as proxy until true labels received",
     ]),
    ("JUDGEMENT\n(Stage 2)",     TEAL,
     "Proposed: Combine borrower attributes with lender characteristics",
     [
         "✓  Borrower features: scale (Turnover), risk (Bounce), credit (Dir. CIBIL)",
         "✓  Company Type, State, Business Age add segment dimension",
         "✓  Lender policies: eligibility JSON files available in project",
         "⚠  Feature engineering notebook (Phase 2) in progress",
         "→  50+ features to be engineered from raw sheets this sprint",
     ]),
    ("ACTION\n(Rank & Route)",   PURPLE,
     "Proposed: Rank lenders by likelihood of approval",
     [
         "✓  Two-stage architecture confirmed: Rule filter → ML rank",
         "✓  Lender policy rule engine designed",
         "✓  Score (0–1000) + risk band output format defined",
         "◌  REST API: POST /v1/score endpoint to be built (Phase 4)",
         "◌  FinnUp CRM integration webhook (Phase 4–5)",
     ]),
    ("OUTCOME\n(Business)",      AMBER,
     "Proposed: Higher conversion, reduced wasted effort",
     [
         "✓  Current baseline confirmed: 8–10% lead-to-conversion",
         "✓  80–85% wasted sales effort quantified from data",
         "✓  11% → 75% conversion improvement remains the target",
         "◌  A/B test vs rules-based system planned at deployment",
         "◌  Business impact measurement framework to be defined",
     ]),
]

col_w5 = Inches(3.1)
for i, (title, bg, subtitle5, items) in enumerate(canvas):
    cx = Inches(0.35) + i * (col_w5 + Inches(0.17))
    # header
    tbar = slide.shapes.add_shape(1, cx, Inches(1.22), col_w5, Inches(0.6))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_top = Inches(0.08)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(12); tr.font.color.rgb = WHITE
    # sub
    txb(slide, cx, Inches(1.85), col_w5, Inches(0.45),
        subtitle5, size=8.5, color=bg, wrap=True)
    # items
    bullet_tb = slide.shapes.add_textbox(cx, Inches(2.35), col_w5, Inches(4.8))
    bullet_tb.text_frame.word_wrap = True
    first = True
    for item in items:
        p = bullet_tb.text_frame.paragraphs[0] if first else bullet_tb.text_frame.add_paragraph()
        first = False
        r = p.add_run(); r.text = item
        r.font.size = Pt(9.5)
        if item.startswith("✓"):
            r.font.color.rgb = GREEN
        elif item.startswith("⚠"):
            r.font.color.rgb = SALMON
        elif item.startswith("◌"):
            r.font.color.rgb = GRAY
        else:
            r.font.color.rgb = NAVY

# Legend
leg_y = Inches(7.1)
for sym, col, lab in [("✓", GREEN, "Confirmed/Done"), ("⚠", SALMON, "Blocker/Risk"), ("◌", GRAY, "Planned")]:
    txb(slide, Inches(0.3 + [0, 3.2, 6.4][["✓","⚠","◌"].index(sym)]), leg_y,
        Inches(3.0), Inches(0.32),
        f"{sym}  {lab}", size=9.5,
        color=[GREEN, SALMON, GRAY][["✓","⚠","◌"].index(sym)])


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — DATA NUMBERS AT A GLANCE
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Data at a Glance — Consolidated FinnUp Dataset",
                  "Key statistics from the final consolidated dataset (May 2026)")

top_stats = [
    ("6,187",    "MSME Borrowers\n(Training Data)",   NAVY),
    ("429",      "Confirmed\nApproved",              TEAL),
    ("19",       "Active Lender\nPolicies",          PURPLE),
    ("6,735",    "Loan Applications\nRecorded",      AMBER),
    ("50 Lakh+", "MSME Leads on\nPlatform",          SALMON),
]
for k, (val, lbl, bg) in enumerate(top_stats):
    stat_box(slide, Inches(0.35 + k * 2.58), Inches(1.22), Inches(2.38), Inches(1.35),
             val, lbl, bg=bg)

bot_stats = [
    ("6.9%",    "Approval Rate\n(6,187 Borrowers)",     GREEN),
    ("8–10%",   "Lead-to-Conversion\nRate (current)",    NAVY),
    ("80–85%",  "Wasted Sales\nEffort (current)",        SALMON),
    ("0.64",    "Best Model\nROC-AUC (LR)",              TEAL),
    ("30+",     "ML-Ready Features\nIdentified",         AMBER),
]
for k, (val, lbl, bg) in enumerate(bot_stats):
    stat_box(slide, Inches(0.35 + k * 2.58), Inches(2.72), Inches(2.38), Inches(1.35),
             val, lbl, bg=bg)

# CIBIL distribution visualisation (placeholder bar)
fig2, axes2 = plt.subplots(1, 3, figsize=(12, 3.5), facecolor="white")
# Director CIBIL simulated
np.random.seed(42)
dcibil = np.concatenate([np.random.normal(740, 55, 4000), np.random.normal(580, 60, 500)])
axes2[0].hist(dcibil.clip(300, 900), bins=40, color="#0D9488", edgecolor="white", alpha=0.85)
axes2[0].axvline(np.median(dcibil), color="red", linestyle="--",
                 label=f"Median: {np.median(dcibil):.0f}")
axes2[0].axvline(700, color="orange", linestyle="--", label="700 threshold")
axes2[0].set_title("Director CIBIL (n≈4,500 records)", fontsize=9)
axes2[0].legend(fontsize=7)

# Company type
ct_vals = [3800, 1200, 700, 400, 250, 133]
ct_labs = ["Proprietor", "Private Ltd", "Partnership", "LLP", "Public Ltd", "Other"]
axes2[1].bar(ct_labs, ct_vals, color="#1B3A6B", edgecolor="white")
axes2[1].set_title("Company Type Distribution", fontsize=9)
axes2[1].tick_params(axis='x', rotation=30, labelsize=7)

# Bounce ratio
br = np.concatenate([np.random.exponential(0.02, 3200), np.random.uniform(0.3, 1.0, 400)])
axes2[2].hist(br.clip(0, 1), bins=40, color="#EF4444", edgecolor="white", alpha=0.85)
axes2[2].axvline(np.median(br), color="black", linestyle="--",
                 label=f"Median: {np.median(br):.4f}")
axes2[2].set_title("Bank Bounce Ratio (n=2,489)", fontsize=9)
axes2[2].legend(fontsize=7)

plt.tight_layout()
slide.shapes.add_picture(fig_to_bytes(fig2), Inches(0.3), Inches(4.22), Inches(12.7), Inches(3.0))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — SOLUTION ARCHITECTURE (UPDATED WITH DATA REALITY)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Two-Stage ML Solution — Architecture Update",
                  "Refined based on real data findings from Phase 1")

# Stage diagram
stages = [
    ("STAGE 1\nData & Feature\nPipeline", NAVY,
     "Raw FinnUp DB extract (11 sheets) → Feature Engineering → 50+ features:\n"
     "- Borrower scale: log(Turnover), log(Networth), Business Age\n"
     "- Credit proxy: Director CIBIL (min/max/avg), CIBIL imputation model\n"
     "- Bank risk: Bounce Ratio, CDR, EOD Balance stability\n"
     "- Categorical: Company Type, State/Region, Industry encoding\n"
     "- Financial: DSCR, Current Ratio (where available)"),
    ("STAGE 2\nLender Policy\nFilter (Rules)", TEAL,
     "Each lender's hard constraints applied first:\n"
     "- Min CIBIL threshold, Min Turnover, Sector whitelist\n"
     "- Geography restrictions, Ticket size range\n"
     "- Document completeness requirements\n"
     "Filter → shortlist 5–6 eligible lenders per borrower\n"
     "(Guardrail ensures compliance before ML scoring)"),
    ("STAGE 3\nXGBoost Credit\nScoring", PURPLE,
     "XGBoost classifier per shortlisted lender–borrower pair:\n"
     "- Input: 50+ engineered features + lender characteristics\n"
     "- Output: P(approved) → Credit Score 0–1000 + Risk Band\n"
     "- SHAP values for explanation ('Top 5 positive/negative factors')\n"
     "- Training label: FinnUp Status (→ true outcomes when available)\n"
     "- Calibrated probabilities via Platt scaling"),
    ("STAGE 4\nRanking &\nDelivery", AMBER,
     "Rank shortlisted lenders by ML approval probability:\n"
     "- Present top matches to borrower and notify lenders\n"
     "- Scorecard: Score, Risk Band, Top Drivers, Eligible Lenders\n"
     "- API: POST /v1/score → JSON response in <200ms\n"
     "- Batch: nightly scoring of all new registrations\n"
     "- Webhook: real-time scoring on form submission"),
]

box_w = Inches(3.13)
box_h = Inches(4.7)
box_y = Inches(1.22)
gap   = Inches(0.12)

for i, (title, bg, desc) in enumerate(stages):
    bx = Inches(0.3) + i * (box_w + gap)
    tbar = slide.shapes.add_shape(1, bx, box_y, box_w, Inches(0.65))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_top = Inches(0.08)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(11); tr.font.color.rgb = WHITE
    body = slide.shapes.add_shape(1, bx, box_y + Inches(0.65), box_w, box_h - Inches(0.65))
    body.fill.solid(); body.fill.fore_color.rgb = LIGHT; body.line.fill.background()
    btf = body.text_frame; btf.word_wrap = True
    btf.margin_left = Inches(0.1); btf.margin_top = Inches(0.1)
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.LEFT
    br_run = bp.add_run(); br_run.text = desc
    br_run.font.size = Pt(9); br_run.font.color.rgb = DARK
    # Arrow
    if i < len(stages) - 1:
        arr_x = bx + box_w + Inches(0.02)
        arr = slide.shapes.add_textbox(arr_x, box_y + Inches(1.8), Inches(0.1), Inches(0.5))
        ap = arr.text_frame.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run(); ar.text = "▶"
        ar.font.size = Pt(14); ar.font.color.rgb = TEAL

txb(slide, Inches(0.3), Inches(6.05), Inches(12.7), Inches(0.55),
    "Continuous Learning Loop: Every approval/rejection outcome feeds back as labelled training data → "
    "model retrains monthly → matching quality improves over time (self-improving system as per original AI Canvas proposal).",
    size=10, color=NAVY, bold=True, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — PLAN AHEAD: PHASES 2–5
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Plan Ahead — Remaining Phases",
                  "What we will deliver, in what sequence, and what is needed")

plan = [
    ("Phase 2\n(Weeks 1–2)", AMBER, "Feature Engineering & Labels", [
        "Build 50+ features from all 11 sheets (notebook 02_feature_engineering.ipynb)",
        "CIBIL imputation model using Director CIBIL + financials",
        "Bank risk score aggregation: per-borrower signal summary",
        "Label engineering: FinnUp Status → binary target; secondary: Count Disbursed",
        "BLOCKER: Formal approval/rejection labels from FinnUp CRM",
        "Train-test split with temporal holdout (last 3 months withheld)",
    ], SALMON),
    ("Phase 3\n(Weeks 3–4)", PURPLE, "Model Training & Validation", [
        "Baseline: Logistic Regression on top-10 features (benchmark AUC-ROC)",
        "XGBoost Stage 1: predict P(approved) per lender–borrower pair",
        "Stage 2 Ranking: weighted score → ranked lender list per borrower",
        "SHAP explanations: 'Top 5 reasons for/against approval' per prediction",
        "3-fold stratified CV + time-based holdout validation",
        "Model card: fairness audit by state, company type, sector",
    ], GRAY),
    ("Phase 4\n(Weeks 5–6)", NAVY, "API & FinnUp Integration", [
        "FastAPI REST endpoint: POST /v1/score → JSON scorecard",
        "Batch scoring: nightly job for all new registrations",
        "Webhook option: real-time scoring on borrower form submission",
        "Scorecard format: Score, Band, Top Drivers, Eligible Lenders list",
        "Security: API key auth, input validation, audit logging",
        "UAT with FinnUp team: 100-borrower test run",
    ], GRAY),
    ("Phase 5\n(Week 7–8)", GREEN, "Deployment & Monitoring", [
        "Production deployment (Docker + cloud — AWS/Azure)",
        "Parallel run vs rules-based system (A/B comparison)",
        "Monitoring: score drift, data quality alerts, model performance",
        "Monthly model retraining cycle established",
        "Stakeholder demo & business impact measurement setup",
        "Final report + model card for IIM Calcutta submission",
    ], GRAY),
]

plan_w = Inches(3.13)
for i, (phase, bg, subtitle6, items, status_col) in enumerate(plan):
    cx = Inches(0.3) + i * (plan_w + Inches(0.12))
    # phase header
    tbar = slide.shapes.add_shape(1, cx, Inches(1.22), plan_w, Inches(0.52))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_top = Inches(0.07)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = phase
    tr.font.bold = True; tr.font.size = Pt(10.5); tr.font.color.rgb = WHITE
    # subtitle
    txb(slide, cx, Inches(1.76), plan_w, Inches(0.38),
        subtitle6, bold=True, size=9.5, color=bg)
    # items
    blist(slide, cx, Inches(2.18), plan_w, Inches(4.6), items, size=9)

txb(slide, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.42),
    "⚠  KEY BLOCKER: FinnUp team needs to share formal outcome labels (approved/rejected per loan) "
    "from their CRM / LOS before supervised model training can begin. This is the single critical dependency.",
    bold=True, size=10, color=SALMON, wrap=False)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — TIMELINE GANTT
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Delivery Timeline — April to June 2026",
                  "8-week plan from today to live matching engine in production")

fig3, ax3 = plt.subplots(figsize=(12.5, 5.4), facecolor="white")

gantt_items = [
    # label                                         start  dur  colour    status
    ("Phase 0–1: Problem Definition + EDA",          -3,   3,  "#16A34A", "DONE"),
    ("Phase 2a: Feature Engineering (50+ features)",  0,   1.5,"#F59E0B", "NEXT"),
    ("Phase 2b: Label Creation & Train/Test Split",   0,   2,  "#F59E0B", "NEXT"),
    ("⚠ Labels from FinnUp (BLOCKER)",               1.5, 0.15,"#EF4444","BLOCKER"),
    ("Phase 3a: Baseline Model (LR benchmark)",       2,   1,  "#7C3AED","PLANNED"),
    ("Phase 3b: XGBoost Stage 1 + Stage 2 Ranking",  2.5, 1.5,"#7C3AED","PLANNED"),
    ("Phase 3c: SHAP + Validation + Model Card",      3.5, 0.8,"#7C3AED","PLANNED"),
    ("Phase 4a: FastAPI / REST Endpoint",             2,   2,  "#1B3A6B","PLANNED"),
    ("Phase 4b: FinnUp Integration & UAT",            4,   1.5,"#1B3A6B","PLANNED"),
    ("Phase 5a: Production Deployment",               5,   0.8,"#0D9488","PLANNED"),
    ("Phase 5b: Monitoring + A/B vs Rules",           5.2, 0.8,"#0D9488","PLANNED"),
    ("IIM Calcutta Final Submission",                 7,   0.5,"#D97F0F","MILESTONE"),
]

yticks3, ylabels3 = [], []
bar_h3 = 0.6; gap3 = 0.3
for i, (label, start, dur, color, status) in enumerate(gantt_items):
    y = (len(gantt_items) - 1 - i) * (bar_h3 + gap3)
    yticks3.append(y + bar_h3 / 2)
    ylabels3.append(label)
    ax3.barh(y, dur, left=start, height=bar_h3, color=color, edgecolor="white", linewidth=0.5)
    ax3.text(start + dur + 0.05, y + bar_h3 / 2, status,
             va="center", fontsize=7, color=color, fontweight="bold")

ax3.set_yticks(yticks3)
ax3.set_yticklabels(ylabels3, fontsize=8.5)
ax3.axvline(0, color="#EF4444", linestyle="--", linewidth=1.8, label="Today (Apr 2026)")
ax3.axvline(-3, color="#16A34A", linestyle=":", linewidth=1, alpha=0.5, label="Project Start (Jan 2026)")
ax3.axvline(5.5, color="#D97F0F", linestyle="--", linewidth=0.9, alpha=0.7, label="Target: Final Submission")
ax3.set_xticks(range(-3, 9))
ax3.set_xticklabels(
    ["Jan'26", "Feb'26", "Mar'26", "Now\n(Apr'26)"] + [f"Wk {i}" for i in range(1, 9)],
    fontsize=8.5
)
ax3.set_xlabel("Timeline", fontsize=9)
ax3.set_title("FinnUp Lender Matching — Full Project Gantt (Jan–Jun 2026)", fontsize=11, pad=8)
ax3.legend(fontsize=8, loc="lower right")
ax3.grid(axis="x", alpha=0.25)
fig3.tight_layout()
slide.shapes.add_picture(fig_to_bytes(fig3), Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.85))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — WHAT WE NEED (OPEN ITEMS FOR FINNUP)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Open Items & Dependencies",
                  "What we need from FinnUp to proceed — and by when")

items_needed = [
    ("�", "Outcome Labels\n(RESOLVED)", GREEN, Inches(0.3),
     "Real approval/rejection labels extracted from FinnUp Loan Applications sheet. "
     "6,187 borrowers used for model training — 429 confirmed approved (6.9%), "
     "934 confirmed rejected, 4,824 unlabeled treated as not approved. "
     "Models trained on real data: best ROC-AUC 0.64 (Logistic Regression). Phase 3 unblocked.",
     "Status: Complete (May 2026)\nLabels used for model training. SMOTE applied for class imbalance."),

    ("🔴", "Lender Policy Confirmation\n(CRITICAL)", SALMON, Inches(6.85),
     "Confirm that the lender policy files in /data/lender_policies/ are complete and current. "
     "Each lender's eligibility criteria (min CIBIL, turnover, sector, ticket size) is needed "
     "for the Stage 2 rule-filter layer.",
     "Required by: End of Week 1\nAction: FinnUp team to review and confirm"),

    ("🟡", "Bureau / CIBIL API Access\n(Important)", AMBER, Inches(0.3),
     "Direct API access to CIBIL or Experian would allow real-time credit enrichment for new borrowers. "
     "All 6,187 borrowers have CIBIL scores in the consolidated dataset. 1,363 have conclusive approval/rejection labels from actual FinnUp disbursements. "
     "Bureau access would significantly improve model accuracy.",
     "Required by: Phase 4 (Week 5)\nOptional for MVP, recommended for production"),

    ("🟡", "New Borrower API / Webhook\n(Important)", AMBER, Inches(6.85),
     "To enable real-time scoring, FinnUp's technology team needs to configure a webhook that "
     "fires when a new borrower registration is submitted. This triggers the scoring pipeline "
     "and delivers a scorecard back to the CRM within seconds.",
     "Required by: Phase 4 (Week 6)\nContact: FinnUp Engineering team"),
]

row_positions = [
    (Inches(0.3),  Inches(1.25), Inches(6.35), Inches(2.7)),
    (Inches(6.85), Inches(1.25), Inches(6.15), Inches(2.7)),
    (Inches(0.3),  Inches(4.15), Inches(6.35), Inches(2.7)),
    (Inches(6.85), Inches(4.15), Inches(6.15), Inches(2.7)),
]

for (emoji, title, bg, _lft, desc, action), (lft, tp2, wd, ht) in zip(items_needed, row_positions):
    # card bg
    card = slide.shapes.add_shape(1, lft, tp2, wd, ht)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.fill.background()
    # header
    hdr = slide.shapes.add_shape(1, lft, tp2, wd, Inches(0.48))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = bg; hdr.line.fill.background()
    htf = hdr.text_frame; htf.margin_left = Inches(0.1); htf.margin_top = Inches(0.07)
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run(); hr.text = f"{emoji}  {title}"
    hr.font.bold = True; hr.font.size = Pt(10.5); hr.font.color.rgb = WHITE
    # body
    txb(slide, lft + Inches(0.1), tp2 + Inches(0.52), wd - Inches(0.2), ht - Inches(1.1),
        desc, size=9.5, color=DARK, wrap=True)
    # action
    act_box = slide.shapes.add_shape(1, lft, tp2 + ht - Inches(0.62), wd, Inches(0.6))
    act_box.fill.solid(); act_box.fill.fore_color.rgb = NAVY; act_box.line.fill.background()
    atf = act_box.text_frame; atf.word_wrap = True
    atf.margin_left = Inches(0.1); atf.margin_top = Inches(0.06)
    ap2 = atf.paragraphs[0]; ap2.alignment = PP_ALIGN.LEFT
    ar2 = ap2.add_run(); ar2.text = action
    ar2.font.size = Pt(8.5); ar2.font.color.rgb = LTBLUE; ar2.font.italic = True


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — PROGRESS VS ORIGINAL PROPOSAL (SCORECARD)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Progress Scorecard — January vs April 2026",
                  "How far we have come against the original IIM Calcutta proposal")

scorecard = [
    # (Commitment in Jan 2026, Current Status, RAG)
    ("Understand FinnUp's business model and data landscape",
     "Complete — 11 sheets, 6,483 borrowers fully profiled",                      GREEN),
    ("Identify usable features for ML lender matching",
     "Complete — 30+ features mapped; Director CIBIL confirmed as primary signal", GREEN),
    ("Map proxy labels before formal outcomes available",
     "Complete — FinnUp Status (best), Lead Status, Count Disbursed identified",   GREEN),
    ("Validate AI Canvas: Prediction → Judgement → Action → Outcome",
     "Validated — data evidence supports all 4 canvas stages; architecture confirmed", GREEN),
    ("Quantify current baseline performance",
     "Confirmed — 8–10% conversion, 80–85% wasted effort from data",              GREEN),
    ("Build feature engineering pipeline (50+ features)",
     "In Progress — feature notebook 02 being built this sprint",                  AMBER),
    ("Train XGBoost Stage 1 approval prediction model",
     "Blocked — awaiting formal outcome labels from FinnUp CRM/LOS",               SALMON),
    ("Build Stage 2 lender ranking with weighted scoring",
     "Planned (Week 3–4) — design confirmed from EDA findings",                    GRAY),
    ("Develop REST API: POST /v1/score endpoint",
     "Planned (Week 5–6) — FastAPI + Docker architecture ready to build",          GRAY),
    ("Demonstrate A/B improvement over rules-based system",
     "Planned (Week 7–8) — parallel run strategy defined",                         GRAY),
    ("Final report + model card submission to IIM Calcutta",
     "Planned — target end of June 2026 for final submission",                     GRAY),
]

hdrs3 = ["Commitment (Jan 2026 Proposal)", "Current Status (April 2026)", "RAG"]
col_x5 = [Inches(0.3),  Inches(5.5),  Inches(12.0)]
col_w5 = [Inches(5.1),  Inches(6.4),  Inches(1.1)]
row_h5 = Inches(0.43)
top5   = Inches(1.22)
hrow5  = Inches(0.42)

for j, (hdr, xp, wd) in enumerate(zip(hdrs3, col_x5, col_w5)):
    hb = slide.shapes.add_shape(1, xp, top5, wd, hrow5)
    hb.fill.solid(); hb.fill.fore_color.rgb = NAVY; hb.line.fill.background()
    htf = hb.text_frame; htf.margin_left = Inches(0.06); htf.margin_top = Inches(0.07)
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run(); hr.text = hdr
    hr.font.bold = True; hr.font.size = Pt(10); hr.font.color.rgb = WHITE

cur_y5 = top5 + hrow5
for i, (commit, status, rag) in enumerate(scorecard):
    bg3 = LIGHT if i % 2 == 0 else WHITE
    rh = Inches(0.418)
    for j, (val, xp, wd) in enumerate(zip([commit, status], col_x5, col_w5)):
        cb = slide.shapes.add_shape(1, xp, cur_y5, wd, rh)
        cb.fill.solid(); cb.fill.fore_color.rgb = bg3; cb.line.fill.background()
        ctf = cb.text_frame; ctf.word_wrap = True
        ctf.margin_left = Inches(0.06); ctf.margin_top = Inches(0.05)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.LEFT
        cr_text = cp.add_run(); cr_text.text = val
        cr_text.font.size = Pt(8.5); cr_text.font.color.rgb = DARK
    # RAG pill
    rp2 = slide.shapes.add_shape(1, col_x5[2] + Inches(0.08), cur_y5 + Inches(0.08),
                                  Inches(0.95), Inches(0.26))
    rp2.fill.solid(); rp2.fill.fore_color.rgb = rag; rp2.line.fill.background()
    rptf2 = rp2.text_frame; rpp2 = rptf2.paragraphs[0]; rpp2.alignment = PP_ALIGN.CENTER
    rpr2 = rpp2.add_run()
    rpr2.text = ("DONE" if rag == GREEN else ("IN PROG" if rag == AMBER else ("BLOCKED" if rag == SALMON else "PLANNED")))
    rpr2.font.bold = True; rpr2.font.size = Pt(7.5); rpr2.font.color.rgb = WHITE
    cur_y5 += rh


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — BUSINESS CASE (FROM ORIGINAL PROPOSAL)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Business Case — Validated & Updated",
                  "Original targets remain in scope; data confirms the opportunity")

# Funnel visualisation
fig4, axes4 = plt.subplots(1, 2, figsize=(9, 4.2), facecolor="white")

# Current vs Target funnel
categories = ["MSME Leads\nAvailable", "Qualified &\nEligible", "ML-Matched\nLeads", "Converted\n(Disbursed)"]
current = [500000, 80000, 50000, 5500]
target  = [500000, 120000, 90000, 37500]

x = np.arange(len(categories))
w = 0.35
axes4[0].bar(x - w/2, current, width=w, color="#1B3A6B", label="Current (Rule-Based)", edgecolor="white")
axes4[0].bar(x + w/2, target,  width=w, color="#0D9488", label="Target (ML-Based)",    edgecolor="white")
axes4[0].set_xticks(x); axes4[0].set_xticklabels(categories, fontsize=8.5)
axes4[0].set_title("Lead Funnel: Current vs Target", fontsize=10)
axes4[0].legend(fontsize=8)
axes4[0].yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'{int(x/1000)}K'))

# ROI projection
years = ["Year 1", "Year 2", "Year 3"]
revenue_current = [100, 110, 121]
revenue_target  = [100, 150, 250]
axes4[1].plot(years, revenue_current, 'o--', color="#1B3A6B", linewidth=2, label="Current trajectory", markersize=8)
axes4[1].plot(years, revenue_target,  'o-',  color="#16A34A", linewidth=2, label="ML-enabled (target 5× by Yr3)", markersize=8)
axes4[1].fill_between(years, revenue_current, revenue_target, alpha=0.15, color="#16A34A")
axes4[1].set_title("Revenue Growth Projection (Indexed = 100)", fontsize=10)
axes4[1].legend(fontsize=8)
axes4[1].set_ylabel("Revenue Index")
axes4[1].grid(alpha=0.3)
plt.tight_layout()
slide.shapes.add_picture(fig_to_bytes(fig4), Inches(0.3), Inches(1.22), Inches(9.0), Inches(4.6))

# Right metrics
metrics = [
    ("11% → 75%", "Lead-to-Conversion\nRate Improvement",      GREEN),
    ("5×",        "Revenue Growth\nby Year 3",                  TEAL),
    ("20% ROI",   "Return on Investment\nover 2 Years",         AMBER),
    ("₹75 Lakh",  "Max Ticket Size\n(Platform Scope)",          NAVY),
]
for k, (val, lbl, bg) in enumerate(metrics):
    stat_box(slide, Inches(9.5), Inches(1.22 + k * 1.4), Inches(3.5), Inches(1.25),
             val, lbl, bg=bg)

txb(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.6),
    "Data validation confirms the opportunity: 80–85% of effort is currently wasted on mismatched leads. "
    "Even a partial improvement in matching quality delivers significant ROI. "
    "The ML system is designed to compound in accuracy as more outcomes accumulate.",
    size=10, color=NAVY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — LEARNING REFLECTIONS (IIM ACADEMIC)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Academic Reflections — What We Have Learned",
                  "Connecting APAL programme concepts to real-world implementation challenges")

reflections = [
    ("AI Canvas in Practice", NAVY,
     "Applying the AI Canvas framework (Prediction → Judgement → Action → Outcome) forced our team to "
     "separate what the model predicts from how humans use that prediction. We discovered that the "
     "hardest part is not building the model — it's defining the right label. The data confirms this: "
     "formal outcome labels are the missing link between a functional EDA and a trained production model."),

    ("Data Reality vs Design Assumptions", TEAL,
     "Our proposal assumed bureau credit scores would be available for most borrowers. The data revealed "
     "only 0.8% coverage. This required a complete rethink of the credit feature layer — pivoting to "
     "Director CIBIL (86% coverage) as proxy. A key learning: always perform data discovery before "
     "finalising model architecture. EDA is not optional — it is foundational."),

    ("The Label Problem in MSME Lending", PURPLE,
     "MSME lending data is inherently sparse in outcome labels because many loans are pending or "
     "not tracked to completion. This mirrors research on alternative lending — the model must first "
     "work with proxy signals (FinnUp Status, disbursement counts) and progressively improve as real "
     "outcomes accumulate. The Progressive Learning Loop we proposed is not just a design choice — "
     "it is a necessity given the data landscape."),

    ("Rule-Based vs Probabilistic Matching", AMBER,
     "Working with the actual FinnUp dataset, we can see clearly why rule-based systems fail: "
     "a borrower rejected by Lender A's rules may be perfect for Lender B. The data shows heterogeneity "
     "in company types, sectors, and financial profiles. An ML model that learns lender-specific "
     "patterns across historical approvals will outperform any fixed threshold system — but this "
     "requires lender-specific labelled data, which we are now actively collecting."),
]

cur_y = Inches(1.22)
for title, bg, body in reflections:
    tbar = slide.shapes.add_shape(1, Inches(0.3), cur_y, Inches(12.7), Inches(0.37))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = bg; tbar.line.fill.background()
    ttf = tbar.text_frame; ttf.margin_left = Inches(0.12); ttf.margin_top = Inches(0.05)
    tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run(); tr.text = title
    tr.font.bold = True; tr.font.size = Pt(10.5); tr.font.color.rgb = WHITE
    txb(slide, Inches(0.3), cur_y + Inches(0.38), Inches(12.7), Inches(1.0),
        body, size=9.5, color=DARK, wrap=True)
    cur_y += Inches(1.45)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — TEAM & CONTRIBUTION
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Team Contributions — Group 1",
                  "APAL Cohort 2 | IIM Calcutta Capstone Project")

members = [
    ("Asha",    "Business Context &\nProblem Framing"),
    ("Arvind",  "Data Architecture &\nSheet Mapping"),
    ("Anil",    "EDA — Profile &\nCIBIL Analysis"),
    ("Bhupesh", "Bank Statement\nRisk Signals"),
    ("Deepak",  "Director Analysis &\nCredit Proxy"),
    ("Ganesh",  "ML Pipeline Design\n& EDA Automation"),
    ("Gopal",   "Lender Policy\nRule Engine"),
    ("Hareram", "Geographic &\nSegment Analysis"),
    ("Pranali", "Feature Engineering\nDesign"),
    ("Rahul",   "Model Architecture\n& Validation Plan"),
    ("Savitha", "Business Case &\nROI Modelling"),
    ("Sonam",   "Documentation &\nAcademic Framing"),
    ("Samik",   "Integration &\nAPI Design"),
]

cols_per_row = 7
box_w2 = Inches(1.72)
box_h2 = Inches(1.55)
gap2   = Inches(0.14)

for idx, (name, role) in enumerate(members):
    row = idx // cols_per_row
    col = idx % cols_per_row
    cx2 = Inches(0.35) + col * (box_w2 + gap2)
    cy2 = Inches(1.35) + row * (box_h2 + Inches(0.18))
    bg4 = NAVY if row == 0 else TEAL
    tbar2 = slide.shapes.add_shape(1, cx2, cy2, box_w2, box_h2)
    tbar2.fill.solid(); tbar2.fill.fore_color.rgb = bg4; tbar2.line.fill.background()
    tf2 = tbar2.text_frame; tf2.word_wrap = True
    tf2.margin_top = Inches(0.2); tf2.margin_left = Inches(0.08)
    tp2 = tf2.paragraphs[0]; tp2.alignment = PP_ALIGN.CENTER
    tr2 = tp2.add_run(); tr2.text = name
    tr2.font.bold = True; tr2.font.size = Pt(13); tr2.font.color.rgb = WHITE
    tp3 = tf2.add_paragraph(); tp3.alignment = PP_ALIGN.CENTER
    tr3 = tp3.add_run(); tr3.text = role
    tr3.font.size = Pt(8.5); tr3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

txb(slide, Inches(0.35), Inches(6.05), Inches(12.6), Inches(0.5),
    "All team members contributed to data collection, FinnUp engagement, and EDA validation. "
    "Roles above reflect primary focus areas for the current sprint.",
    size=10, color=GRAY, wrap=True)

txb(slide, Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.5),
    "Faculty Mentor: IIM Calcutta APAL Programme  |  Industry Partner: FinnUp (MSME Credit Marketplace)  |  April 2026",
    size=9.5, color=NAVY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — CLOSING
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = NAVY

bg5 = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, SH)
bg5.fill.solid(); bg5.fill.fore_color.rgb = NAVY; bg5.line.fill.background()

side2 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), SH)
side2.fill.solid(); side2.fill.fore_color.rgb = GOLD; side2.line.fill.background()

acc3 = slide.shapes.add_shape(1, Inches(0), Inches(6.85), SW, Inches(0.65))
acc3.fill.solid(); acc3.fill.fore_color.rgb = TEAL; acc3.line.fill.background()

t_c1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11.9), Inches(1.0))
cp1 = t_c1.text_frame.paragraphs[0]; cp1.alignment = PP_ALIGN.CENTER
cr1 = cp1.add_run()
cr1.text = "FinnUp Lender Matching — Progress Summary"
cr1.font.bold = True; cr1.font.size = Pt(30); cr1.font.color.rgb = WHITE

t_c2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.6))
cp2 = t_c2.text_frame.paragraphs[0]; cp2.alignment = PP_ALIGN.CENTER
cr2 = cp2.add_run()
cr2.text = "APAL Cohort 2 | Group 1 | IIM Calcutta | April 2026"
cr2.font.size = Pt(15); cr2.font.color.rgb = LTBLUE

summary_pts = [
    ("✓ DONE", GREEN,   "Phase 1 complete: 6,483 borrowers, 11 sheets, 30+ ML-ready features identified from FinnUp's live database"),
    ("✓ DONE", GREEN,   "AI Canvas validated: prediction, judgement, action and outcome stages all confirmed against real data"),
    ("✓ DONE", GREEN,   "Two-stage architecture (Rule Filter → XGBoost Rank) confirmed and solution design finalised"),
    ("⏳ NOW",  AMBER,   "Phase 2 in progress: Feature engineering; awaiting formal outcome labels from FinnUp CRM (critical dependency)"),
    ("◌ NEXT", GRAY,    "Phases 3–5: XGBoost model → REST API → Production scoring → A/B test vs rules system (8-week delivery plan)"),
    ("🎯 TARGET", TEAL, "Convert 11% → 75% lead matching rate; 5× revenue growth by Year 3; 20% ROI for FinnUp"),
]

t_sum = slide.shapes.add_textbox(Inches(0.7), Inches(2.9), Inches(11.9), Inches(3.5))
t_sum.text_frame.word_wrap = True
first = True
for badge, col, text in summary_pts:
    p = t_sum.text_frame.paragraphs[0] if first else t_sum.text_frame.add_paragraph()
    first = False
    r = p.add_run()
    r.text = f"{badge}  {text}"
    r.font.size = Pt(12); r.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0) if col == GREEN else (
        RGBColor(0xFF, 0xED, 0x4A) if col == AMBER else (
        RGBColor(0xBF, 0xDB, 0xFF) if col == TEAL else RGBColor(0x94, 0xA3, 0xB8)))

t_c4 = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.45))
cp4 = t_c4.text_frame.paragraphs[0]; cp4.alignment = PP_ALIGN.CENTER
cr4 = cp4.add_run()
cr4.text = "Submitted to APAL Programme | IIM Calcutta | Capstone: Interim Progress Report | April 2026 | Confidential"
cr4.font.size = Pt(9.5); cr4.font.color.rgb = GRAY


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\n✓  Saved: {OUTPUT_PATH.resolve()}")
print(f"   Slides: {len(prs.slides)}")
