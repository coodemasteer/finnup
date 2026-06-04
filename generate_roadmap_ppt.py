"""
generate_roadmap_ppt.py
-----------------------
Generates  outputs/FinnUp_ProjectStatus_Roadmap.pptx
  – Project Status (what's done) + full roadmap with timeline
    for the FinnUp MSME Credit Intelligence Platform

Run:
    python generate_roadmap_ppt.py
"""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
import warnings
warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Output ────────────────────────────────────────────────────────────────────
OUTPUT_PATH = Path("outputs") / "FinnUp_ProjectStatus_Roadmap.pptx"
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x3A, 0x6B)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
SALMON  = RGBColor(0xEF, 0x44, 0x44)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
GRAY    = RGBColor(0x64, 0x74, 0x8B)
LTBLUE  = RGBColor(0x93, 0xC5, 0xFD)

SW = Inches(13.33)
SH = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def fig_to_bytes(fig, dpi=130):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def add_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    tf.margin_left = Inches(0.35)
    tf.margin_top  = Inches(0.18)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.35), Inches(0.82), Inches(12), Pt(18))
        sp = sub.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(11)
        sr.font.color.rgb = TEAL
        sr.font.italic = True
    line = slide.shapes.add_shape(1, Inches(0), Inches(1.1), SW, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    return slide


def add_textbox(slide, left, top, width, height, text,
                bold=False, size=11, color=DARK,
                align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text  = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, left, top, width, height, items,
                   size=11, title=None, title_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.bold  = True
        r.font.size  = Pt(size + 1)
        r.font.color.rgb = title_color or NAVY
    for item in items:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size  = Pt(size)
        r.font.color.rgb = DARK
    return tb


def add_stat_box(slide, left, top, width, height, value, label, bg=NAVY):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.bold  = True
    r.font.size  = Pt(24)
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size  = Pt(9)
    r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)


def pill(slide, left, top, width, height, text, bg=TEAL, font_size=10):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = text
    r.font.bold  = True
    r.font.size  = Pt(font_size)
    r.font.color.rgb = WHITE


# =============================================================================
# BUILD PRESENTATION
# =============================================================================
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = NAVY

band = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(5.0))
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()

acc = slide.shapes.add_shape(1, Inches(0), Inches(5.0), SW, Pt(6))
acc.fill.solid()
acc.fill.fore_color.rgb = TEAL
acc.line.fill.background()

t1 = slide.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.2))
p1 = t1.text_frame.paragraphs[0]
r1 = p1.add_run()
r1.text = "FinnUp MSME Credit Intelligence Platform"
r1.font.bold  = True
r1.font.size  = Pt(36)
r1.font.color.rgb = WHITE

t2 = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.8))
p2 = t2.text_frame.paragraphs[0]
r2 = p2.add_run()
r2.text = "Project Status Report & Delivery Roadmap"
r2.font.size  = Pt(22)
r2.font.color.rgb = LTBLUE

t3 = slide.shapes.add_textbox(Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.6))
p3 = t3.text_frame.paragraphs[0]
r3 = p3.add_run()
r3.text = "April 2026  |  MSME Marketplace — Lenders × Borrowers  |  Automated Lending Decisioning"
r3.font.size  = Pt(13)
r3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF)

# status chips
chip_data = [("Phase 1: EDA", GREEN), ("Phase 2: Features", AMBER), ("Phase 3: Model", GRAY), ("Phase 4: API", GRAY)]
for i, (txt, col) in enumerate(chip_data):
    cx = Inches(0.9) + i * Inches(3.1)
    sh = slide.shapes.add_shape(1, cx, Inches(4.3), Inches(2.8), Inches(0.42))
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = False
    pp = tf.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run()
    rr.text = ("✓ COMPLETE  " if col == GREEN else ("⏳ NEXT  " if col == AMBER else "◌ PLANNED  ")) + txt
    rr.font.bold  = True
    rr.font.size  = Pt(10)
    rr.font.color.rgb = WHITE

t4 = slide.shapes.add_textbox(Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.8))
t4.text_frame.word_wrap = True
p4 = t4.text_frame.paragraphs[0]
r4 = p4.add_run()
r4.text = (
    "This document captures the completed EDA milestone and the end-to-end delivery plan for an "
    "AI-powered credit scoring engine that plugs into the FinnUp marketplace — "
    "enabling automated lead evaluation for lenders and instant creditworthiness signals for borrowers."
)
r4.font.size  = Pt(12)
r4.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — EXECUTIVE SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Executive Summary",
                  "What has been done · Where we are · What comes next")

# 3-column layout
cols = [
    ("✓ COMPLETED\n(Phase 1 — EDA)", GREEN, [
        "Loaded & audited 6,483 borrowers across 11 data sheets",
        "Identified 30+ usable ML features with >50% coverage",
        "Mapped proxy labels: FinnUp Status, Lead Status, Count Disbursed",
        "Analysed bank risk signals: bounce ratio, EOD balance, CDR",
        "Profiled director CIBIL as key proxy for borrower credit",
        "Documented geographic & industry distribution",
        "Full EDA deck auto-generated (FinnUp_EDA_Report.pptx)",
    ]),
    ("⏳ NEXT SPRINT\n(Phase 2 — Features)", AMBER, [
        "Receive outcome labels from FinnUp team",
        "Engineer 50+ model-ready features from raw sheets",
        "Build CIBIL-score imputation using director proxy",
        "Create fraud & risk flag features from bank data",
        "Encode lender-policy eligibility rules as features",
        "Train/test split with temporal holdout",
        "Deliver feature store & feature importance analysis",
    ]),
    ("◌ PLANNED\n(Phase 3–5: Model → API → Live)", NAVY, [
        "Train XGBoost / LightGBM credit scoring model",
        "Calibrate probability thresholds per lender policy",
        "Build REST API: POST /score with borrower JSON payload",
        "Integrate webhook or batch pipeline into FinnUp platform",
        "Automated lead scoring on new registrations",
        "Monitoring dashboard for score drift & data quality",
        "Optional: real-time bureau pull (CIBIL/Experian) enrichment",
    ]),
]

col_x = [Inches(0.35), Inches(4.55), Inches(8.75)]
for ci, (title, col, items) in enumerate(cols):
    cx = col_x[ci]
    # header box
    hb = slide.shapes.add_shape(1, cx, Inches(1.25), Inches(3.95), Inches(0.55))
    hb.fill.solid()
    hb.fill.fore_color.rgb = col
    hb.line.fill.background()
    htf = hb.text_frame
    htf.word_wrap = False
    htf.margin_left = Inches(0.1)
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run()
    hr.text = title
    hr.font.bold  = True
    hr.font.size  = Pt(10)
    hr.font.color.rgb = WHITE
    # bullet box
    add_bullet_box(slide, cx, Inches(1.82), Inches(3.95), Inches(5.4),
                   items, size=9.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — PHASE 1 EDA COMPLETED
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 1 — EDA: What Was Done",
                  "Comprehensive data audit across all 11 FinnUp data sheets")

# Data coverage bar chart
feature_data = {
    "Turnover":            100.0,
    "Networth":            92.0,
    "Company Name":        96.0,
    "Business Age":        51.0,
    "Director CIBIL":      86.0,
    "Bank Account Data":   38.0,
    "Financial KPIs":      7.7,
    "State / Location":    55.0,
    "Lead Status":         0.5,
    "CIBIL Score":         0.8,
}
labels = list(feature_data.keys())
values = list(feature_data.values())
colors_bar = ["#16A34A" if v >= 80 else ("#F59E0B" if v >= 40 else "#EF4444") for v in values]

fig, ax = plt.subplots(figsize=(7.5, 4.2), facecolor="white")
bars = ax.barh(labels, values, color=colors_bar, edgecolor="white", height=0.6)
ax.axvline(50, color="gray", linestyle="--", linewidth=0.9, label="50% threshold")
for bar, v in zip(bars, values):
    ax.text(v + 0.8, bar.get_y() + bar.get_height() / 2,
            f"{v:.0f}%", va="center", fontsize=9)
ax.set_xlim(0, 112)
ax.set_xlabel("Coverage (%)")
ax.set_title("Feature Data Coverage across 6,483 Borrowers", fontsize=11)
green_p = mpatches.Patch(color="#16A34A", label="≥80% — Production Ready")
amber_p = mpatches.Patch(color="#F59E0B", label="40–80% — Usable with care")
red_p   = mpatches.Patch(color="#EF4444", label="<40% — Sparse / Needs fix")
ax.legend(handles=[green_p, amber_p, red_p], fontsize=8, loc="lower right")
plt.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.25), Inches(7.8), Inches(5.6))

add_bullet_box(slide, Inches(8.35), Inches(1.30), Inches(4.7), Inches(5.7),
               ["11 sheets, 6,483 borrowers loaded & audited",
                "Turnover & Networth: near-complete — primary scale signals",
                "Director CIBIL (86%): best proxy for borrower creditworthiness",
                "Bank data covers 38% of borrowers — strong risk signals",
                "Borrower CIBIL: only 0.8% — imputation needed",
                "Financial KPIs: 500 borrowers (7.7%) with formal financials",
                "Lead Status: 99%+ unlabelled — outcome labels awaited",
                "FinnUp Status (Loans sheet): viable binary training label",
                "Geographic coverage: 55% have state data",
                "30+ features identified as model-ready after engineering",
               ], title="EDA Key Findings", size=10)

# DONE badge
done_sh = slide.shapes.add_shape(1, Inches(9.5), Inches(5.75), Inches(2.0), Inches(0.52))
done_sh.fill.solid()
done_sh.fill.fore_color.rgb = GREEN
done_sh.line.fill.background()
dtf = done_sh.text_frame
dp  = dtf.paragraphs[0]
dp.alignment = PP_ALIGN.CENTER
dr  = dp.add_run()
dr.text = "✓  MILESTONE COMPLETE"
dr.font.bold = True
dr.font.size = Pt(10)
dr.font.color.rgb = WHITE


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — EDA KEY NUMBERS
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "EDA Key Numbers at a Glance",
                  "Derived from FinnUp_Borrowers.xlsx — 11 sheets analysed")

stats_top = [
    ("6,483",  "Total Borrowers",      NAVY),
    ("11",     "Data Sheets",          TEAL),
    ("9,273",  "Director Records",     NAVY),
    ("3,996",  "Bank Statement Rows",  TEAL),
    ("2,489",  "Borrowers with Bank\nAccount Data", NAVY),
]
for k, (val, lbl, bg) in enumerate(stats_top):
    sx = Inches(0.35) + k * Inches(2.59)
    add_stat_box(slide, sx, Inches(1.25), Inches(2.4), Inches(1.35), val, lbl, bg=bg)

stats_bot = [
    ("752",  "Median CIBIL Score\n(borrower, sparse)",      AMBER),
    ("38%",  "Bank Data Coverage\n(2,489 / 6,483)",         PURPLE),
    ("6%",   "Borrowers with any\nDisbursement (Count>0)",  SALMON),
    ("30+",  "Features Identified\nfor ML",                 GREEN),
    ("500",  "Borrowers with\nFormal Financial KPIs",       TEAL),
]
for k, (val, lbl, bg) in enumerate(stats_bot):
    sx = Inches(0.35) + k * Inches(2.59)
    add_stat_box(slide, sx, Inches(2.72), Inches(2.4), Inches(1.35), val, lbl, bg=bg)

# Second block — feature groups breakdown table
features_table = [
    ("PROFILE SHEET",   "Turnover, Networth, Business Age, Company Type, State, Contact",  "100% / 92%",  GREEN),
    ("DIRECTORS",       "Director CIBIL, Ownership %, Number of Directors",                "86%",         GREEN),
    ("BANK STATEMENTS", "Bounce Ratio, EOD Balance, Credit/Debit Ratio, Cheque Bounces",   "38%",         AMBER),
    ("LOAN REQUESTS",   "Product Name, Loan Amount, FinnUp Status (label)",                "67%",         AMBER),
    ("FINANCIAL KPIs",  "DSCR, Current Ratio, Net Profit Margin, Debt/Equity",             "7.7%",        SALMON),
    ("DOCUMENTS",       "Document type completeness",                                      "67%",         AMBER),
]
hdrs = ["Feature Source", "Key Features", "Coverage", "Readiness"]
col_x2   = [Inches(0.35), Inches(2.65), Inches(9.05), Inches(10.75)]
col_w2   = [Inches(2.2),  Inches(6.3),  Inches(1.6),  Inches(2.2)]
row_h2   = Inches(0.40)
top2     = Inches(4.20)

for j, (hdr, xp, wd) in enumerate(zip(hdrs, col_x2, col_w2)):
    hb = slide.shapes.add_shape(1, xp, top2, wd, row_h2)
    hb.fill.solid()
    hb.fill.fore_color.rgb = NAVY
    hb.line.fill.background()
    htf = hb.text_frame
    htf.margin_left = Inches(0.06)
    htf.margin_top  = Inches(0.08)
    hp  = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr  = hp.add_run()
    hr.text = hdr
    hr.font.bold  = True
    hr.font.size  = Pt(10)
    hr.font.color.rgb = WHITE

for i, (src, feats, cov, clr) in enumerate(features_table):
    y = top2 + row_h2 + i * row_h2
    bg = LIGHT if i % 2 == 0 else WHITE
    for j, (val, xp, wd) in enumerate(zip([src, feats, cov, ""], col_x2, col_w2)):
        cb = slide.shapes.add_shape(1, xp, y, wd, row_h2)
        cb.fill.solid()
        cb.fill.fore_color.rgb = bg
        cb.line.fill.background()
        ctf = cb.text_frame
        ctf.margin_left = Inches(0.06)
        ctf.margin_top  = Inches(0.06)
        cp  = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.LEFT
        cr  = cp.add_run()
        cr.text = val
        cr.font.size  = Pt(9)
        cr.font.color.rgb = DARK
    # Readiness pill
    rp = slide.shapes.add_shape(1, col_x2[3] + Inches(0.1), y + Inches(0.07),
                                 Inches(1.8), Inches(0.26))
    rp.fill.solid()
    rp.fill.fore_color.rgb = clr
    rp.line.fill.background()
    rtf = rp.text_frame
    rp2  = rtf.paragraphs[0]
    rp2.alignment = PP_ALIGN.CENTER
    rr2 = rp2.add_run()
    rr2.text = ("READY" if clr == GREEN else ("PARTIAL" if clr == AMBER else "SPARSE"))
    rr2.font.bold  = True
    rr2.font.size  = Pt(8)
    rr2.font.color.rgb = WHITE


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — END-TO-END SOLUTION VISION
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "End-to-End Solution Vision",
                  "FinnUp Credit Intelligence: from lead registration to lender match")

add_textbox(slide, Inches(0.35), Inches(1.2), Inches(12.6), Inches(0.4),
            "The goal: any MSME borrower registering on the FinnUp platform receives an instant "
            "credit score & lender eligibility match — and lenders receive pre-scored, ranked leads.",
            size=11, color=GRAY, wrap=True)

# Pipeline boxes
pipeline = [
    ("1\nLEAD\nINGESTION",   "Borrower registers on FinnUp portal / CRM / Excel upload",   NAVY),
    ("2\nDATA\nENRICHMENT",  "Auto-pull bank data, director info, bureau scores (optional)", TEAL),
    ("3\nFEATURE\nENGINEERING","Compute 50+ signals: bounce ratio, DSCR, CIBIL proxy, etc.", PURPLE),
    ("4\nML CREDIT\nSCORING", "XGBoost model outputs credit score (0–1000) + risk band",    AMBER),
    ("5\nLENDER\nMATCHING",  "Score filtered against each lender's policy rules → ranked list", SALMON),
    ("6\nOUTPUT\n& ACTION",  "Dashboard alert / API JSON / Email to lender with scorecard",  GREEN),
]

box_w = Inches(1.9)
box_h = Inches(2.3)
box_y = Inches(1.9)
gap   = Inches(0.25)
start = Inches(0.35)

for i, (lbl, desc, bg) in enumerate(pipeline):
    bx = start + i * (box_w + gap)
    # Box
    sh = slide.shapes.add_shape(1, bx, box_y, box_w, box_h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.1)
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r  = p.add_run()
    r.text = lbl
    r.font.bold  = True
    r.font.size  = Pt(11)
    r.font.color.rgb = WHITE
    # Arrow
    if i < len(pipeline) - 1:
        arr_x = bx + box_w + Inches(0.05)
        arr = slide.shapes.add_textbox(arr_x, box_y + Inches(0.9), gap + Inches(0.05), Inches(0.5))
        ap  = arr.text_frame.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar  = ap.add_run()
        ar.text = "▶"
        ar.font.size  = Pt(14)
        ar.font.color.rgb = GRAY
    # Description below box
    desc_tb = slide.shapes.add_textbox(bx, box_y + box_h + Inches(0.12),
                                        box_w, Inches(1.3))
    dp = desc_tb.text_frame
    dp.word_wrap = True
    dpp = dp.paragraphs[0]
    dpp.alignment = PP_ALIGN.CENTER
    dr  = dpp.add_run()
    dr.text = desc
    dr.font.size  = Pt(8.5)
    dr.font.color.rgb = DARK

# Integration options note
add_textbox(slide, Inches(0.35), Inches(6.3), Inches(12.6), Inches(0.8),
            "Integration options: (a) FinnUp embeds scoring API in their portal  "
            "│  (b) Batch job scores new registrations nightly  "
            "│  (c) Webhook triggered on borrower form-submit",
            size=10, color=TEAL, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — PLATFORM INTEGRATION ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Platform Integration Architecture",
                  "How the scoring engine plugs into the FinnUp ecosystem")

arch_items = [
    # Left column: FinnUp Platform
    ("FinnUp Platform", NAVY, [
        "Borrower Registration Portal",
        "Admin / CRM Dashboard",
        "Lender Onboarding Module",
        "Excel / CSV Bulk Upload",
        "Existing Borrower Database",
    ], Inches(0.3), Inches(1.25), Inches(3.2), Inches(5.5)),
    # Middle column: Integration Layer
    ("Integration Layer\n(NEW)", TEAL, [
        "REST API  POST /v1/score",
        "Batch Scoring Job (nightly cron)",
        "Webhook Listener (real-time)",
        "Auth: API Key + mTLS",
        "Rate-limited & audit-logged",
    ], Inches(3.85), Inches(1.25), Inches(3.2), Inches(5.5)),
    # Right column: ML Backend
    ("ML Scoring Engine\n(NEW)", PURPLE, [
        "Feature Engineering Pipeline",
        "XGBoost Credit Score Model",
        "Lender Policy Rule Engine",
        "Score Explanation (SHAP)",
        "Model Monitoring & Drift Alerts",
    ], Inches(7.4), Inches(1.25), Inches(3.2), Inches(5.5)),
    # Far-right: Outputs
    ("Outputs & Consumers", AMBER, [
        "Lender Dashboard: ranked leads",
        "Borrower: pre-qualification score",
        "Alert: email / WhatsApp to lender",
        "Analytics: conversion tracking",
        "Reports: weekly portfolio summary",
    ], Inches(10.95), Inches(1.25), Inches(2.1), Inches(5.5)),
]

for title, bg, items, lft, tp, wd, ht in arch_items:
    # Header
    hb = slide.shapes.add_shape(1, lft, tp, wd, Inches(0.52))
    hb.fill.solid()
    hb.fill.fore_color.rgb = bg
    hb.line.fill.background()
    htf = hb.text_frame
    htf.word_wrap = True
    htf.margin_top = Inches(0.06)
    htf.margin_left = Inches(0.08)
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run()
    hr.text = title
    hr.font.bold = True
    hr.font.size = Pt(10)
    hr.font.color.rgb = WHITE
    # Body
    bb = slide.shapes.add_shape(1, lft, tp + Inches(0.52), wd, ht - Inches(0.52))
    bb.fill.solid()
    bb.fill.fore_color.rgb = LIGHT
    bb.line.fill.background()
    btf = bb.text_frame
    btf.word_wrap = True
    btf.margin_top = Inches(0.12)
    btf.margin_left = Inches(0.1)
    first = True
    for item in items:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(9.5)
        r.font.color.rgb = DARK

# Arrows between columns
arrow_positions = [Inches(3.52), Inches(7.07), Inches(10.62)]
for ax_ in arrow_positions:
    arr = slide.shapes.add_textbox(ax_, Inches(3.65), Inches(0.32), Inches(0.35))
    ap  = arr.text_frame.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    ar  = ap.add_run()
    ar.text = "▶"
    ar.font.size = Pt(16)
    ar.font.color.rgb = TEAL

add_textbox(slide, Inches(0.3), Inches(6.92), Inches(12.7), Inches(0.45),
            "Data flows: FinnUp Platform → Integration Layer (secure API/batch) → ML Engine → Outputs back to FinnUp",
            size=10, color=GRAY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — PHASE 2: FEATURE ENGINEERING (NEXT SPRINT)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 2 — Feature Engineering & Label Creation",
                  "Weeks 1–2 of remaining roadmap | Requires outcome labels from FinnUp")

left_items = [
    ("Profile Features", NAVY, [
        "log(Turnover), log(Networth), Business Age bins",
        "Company Type one-hot encoding",
        "State / Region encoding (grouped)",
        "Borrower-to-director CIBIL ratio",
        "Documents completeness score",
        "Application-to-sanction gap (days)",
    ]),
    ("Bank Statement Features", TEAL, [
        "Bounce Ratio (total & trend)",
        "Avg / Median EOD Balance (log)",
        "Credit-Debit Ratio",
        "Inward vs Outward bounce breakdown",
        "Stability score (σ of monthly balance)",
        "Bank vintage (months of history)",
    ]),
]

right_items = [
    ("Director / Credit Features", PURPLE, [
        "Min / Max / Avg Director CIBIL",
        "Number of directors",
        "Promoter ownership % (max)",
        "CIBIL imputation flag (proxy used?)",
        "Director default flag (CIBIL < 600)",
    ]),
    ("Financial KPI Features (500 borrowers)", AMBER, [
        "DSCR (Avg / Min) — most predictive",
        "Current Ratio, Debt-Equity Ratio",
        "Net Profit Margin, PBDITA Margin",
        "YoY Revenue growth (%)",
        "KPI availability flag (for 93% without)",
    ]),
    ("Label Engineering", SALMON, [
        "Primary: FinnUp Status (YES/NO) per loan request",
        "Secondary: Count Disbursed > 0 binary flag",
        "Tertiary: Lead Status (once populated)",
        "True labels: awaited from FinnUp team → BLOCKER",
    ]),
]

col_w = Inches(3.7)
for ci, col_data in enumerate([left_items, right_items]):
    cx = Inches(0.35) + ci * Inches(6.65)
    cur_y = Inches(1.3)
    for gname, gbg, gitems in col_data:
        # group header
        gh = slide.shapes.add_shape(1, cx, cur_y, col_w * (1 if ci == 0 else 1.7), Inches(0.38))
        gh.fill.solid()
        gh.fill.fore_color.rgb = gbg
        gh.line.fill.background()
        gtf = gh.text_frame
        gtf.margin_left = Inches(0.1)
        gtf.margin_top  = Inches(0.05)
        gp  = gtf.paragraphs[0]
        gp.alignment = PP_ALIGN.LEFT
        gr  = gp.add_run()
        gr.text = gname
        gr.font.bold  = True
        gr.font.size  = Pt(10)
        gr.font.color.rgb = WHITE
        cur_y += Inches(0.4)
        for item in gitems:
            add_textbox(slide, cx + Inches(0.1), cur_y,
                        col_w * (1 if ci == 0 else 1.7), Inches(0.3),
                        f"▸  {item}", size=9, color=DARK)
            cur_y += Inches(0.28)
        cur_y += Inches(0.15)

add_textbox(slide, Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.38),
            "⚠  BLOCKER: True outcome labels from FinnUp are needed before supervised model training can begin.",
            bold=True, size=10, color=SALMON, wrap=False)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — PHASE 3: MODEL DEVELOPMENT
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 3 — Model Development",
                  "Weeks 3–4 | XGBoost credit scoring + lender policy rule engine")

models = [
    ("Credit Approval\nClassifier", NAVY,
     "XGBoost / LightGBM\nbinary classifier",
     "Predicts P(approved) per borrower using all engineered features",
     ["AUC-ROC > 0.80 target", "F1 on minority class", "PSI < 0.2 (stability)"]),

    ("Credit Score\n(0–1000 band)", TEAL,
     "Score = 1000 × calibrated\nprobability",
     "Converts model output to human-readable credit score with band labels",
     ["Score bands: Poor / Fair / Good / Very Good / Excellent",
      "Calibrated using Platt scaling",
      "SHAP explanations per score"]),

    ("Lender Policy\nRule Engine", PURPLE,
     "Rule tree per lender\n(loan amount, sector, CIBIL min)",
     "Filters eligible borrower-lender pairs after scoring",
     ["Each lender has policy JSON with eligibility rules",
      "Scoring first, then eligibility filter",
      "Lender-level pass/fail with reason codes"]),

    ("CIBIL Score\nImputer", AMBER,
     "Regression model:\nDirector CIBIL → Borrower CIBIL",
     "Fills CIBIL gaps (99% missing) using director & financial features",
     ["Trained on the 55 borrowers with actual CIBIL",
      "Flagged as 'imputed' in output",
      "Uncertainty interval provided"]),
]

mw = Inches(3.0)
gap = Inches(0.18)
for i, (title, bg, tech, desc, kpis) in enumerate(models):
    mx = Inches(0.35) + i * (mw + gap)
    # title bar
    tbar = slide.shapes.add_shape(1, mx, Inches(1.25), mw, Inches(0.55))
    tbar.fill.solid()
    tbar.fill.fore_color.rgb = bg
    tbar.line.fill.background()
    ttf = tbar.text_frame
    ttf.word_wrap = False
    ttf.margin_top = Inches(0.08)
    tp  = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr  = tp.add_run()
    tr.text = title
    tr.font.bold = True
    tr.font.size = Pt(11)
    tr.font.color.rgb = WHITE
    # tech box
    tech_b = slide.shapes.add_shape(1, mx, Inches(1.80), mw, Inches(0.55))
    tech_b.fill.solid()
    tech_b.fill.fore_color.rgb = LIGHT
    tech_b.line.fill.background()
    tetf = tech_b.text_frame
    tetf.margin_left = Inches(0.08)
    tetf.margin_top  = Inches(0.06)
    te_p = tetf.paragraphs[0]
    te_r = te_p.add_run()
    te_r.text = tech
    te_r.font.size = Pt(9)
    te_r.font.color.rgb = GRAY
    te_r.font.italic = True
    # description
    add_textbox(slide, mx, Inches(2.42), mw, Inches(0.9),
                desc, size=9.5, color=DARK, wrap=True)
    # KPIs / bullets
    add_bullet_box(slide, mx, Inches(3.35), mw, Inches(2.5),
                   kpis, size=9.5, title="Success Metrics / Details")

add_bullet_box(slide, Inches(0.35), Inches(5.9), Inches(12.6), Inches(1.4),
               ["Baseline: logistic regression with top-10 features → benchmark for XGBoost",
                "Validation: 3-fold stratified CV + time-based holdout (last 3 months of data)",
                "Fairness check: score distributions by state, company type, size segment",
                "Model card & explainability report generated at end of phase",
               ], title="Validation Strategy", size=10)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — PHASE 4: API & INTEGRATION
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 4 — API, Integration & Automation",
                  "Weeks 5–6 | Scoring engine accessible via REST API + batch pipeline")

left_bullets = [
    ("REST API Endpoints", NAVY, [
        "POST  /v1/score          → Single borrower scoring",
        "POST  /v1/score/batch   → Bulk CSV / JSON scoring",
        "GET   /v1/score/{id}    → Retrieve cached score",
        "GET   /v1/lenders/match → Top lenders for a borrower",
        "GET   /v1/health        → Uptime & model version",
    ]),
    ("Security & Reliability", TEAL, [
        "API Key authentication (per lender / per team)",
        "Rate limiting: 100 req/min per key",
        "Input validation & sanitisation on all fields",
        "Audit log: every request logged with timestamp & user",
        "HTTPS enforced; secrets in environment variables",
        "99.5% uptime SLA target (cloud-hosted)",
    ]),
]

right_bullets = [
    ("Automated Lead Pipeline", PURPLE, [
        "Nightly batch job: score all new registrations since last run",
        "Webhook option: trigger score on form submit (sub-second)",
        "FinnUp CRM integration: push score + lender match to lead record",
        "Email / WhatsApp alert to matched lenders with scorecard PDF",
        "Excel upload UI for ops teams to score bulk leads manually",
    ]),
    ("Scorecard Output Format", AMBER, [
        "Borrower ID, Credit Score (0–1000), Risk Band",
        "Top 5 positive / negative feature drivers (SHAP)",
        "Eligible lenders list with product match reasons",
        "Data quality flags (imputed fields, missing docs)",
        "Recommended next action per lender policy",
    ]),
]

for ci, col_data in enumerate([left_bullets, right_bullets]):
    cx  = Inches(0.35) + ci * Inches(6.5)
    cur_y = Inches(1.3)
    for gname, gbg, gitems in col_data:
        gh = slide.shapes.add_shape(1, cx, cur_y, Inches(6.2), Inches(0.38))
        gh.fill.solid()
        gh.fill.fore_color.rgb = gbg
        gh.line.fill.background()
        gtf = gh.text_frame
        gtf.margin_left = Inches(0.1)
        gtf.margin_top  = Inches(0.05)
        gp  = gtf.paragraphs[0]
        gp.alignment = PP_ALIGN.LEFT
        gr  = gp.add_run()
        gr.text = gname
        gr.font.bold  = True
        gr.font.size  = Pt(10)
        gr.font.color.rgb = WHITE
        cur_y += Inches(0.4)
        for item in gitems:
            add_textbox(slide, cx + Inches(0.1), cur_y, Inches(6.1), Inches(0.3),
                        f"▸  {item}", size=9, color=DARK)
            cur_y += Inches(0.29)
        cur_y += Inches(0.18)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — PHASE 5: PRODUCTION & MONITORING
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Phase 5 — Production Deployment & Monitoring",
                  "Weeks 7–8 | Live scoring, dashboards, model refresh cycle")

prod_sections = [
    ("Deployment", GREEN, Inches(0.35), [
        "Cloud deployment: AWS / Azure / GCP (containerised)",
        "Docker image with model artefacts + API server",
        "Auto-scaling: handles FinnUp's growth trajectory",
        "Blue-green deployment for zero-downtime model updates",
        "Infrastructure as Code (Terraform / Bicep) for reproducibility",
    ]),
    ("Monitoring Dashboard", TEAL, Inches(3.55), [
        "Score distribution daily snapshot",
        "Feature drift alerts (PSI > 0.2 threshold)",
        "Model performance vs reported outcomes",
        "Data quality: null rate trends per field",
        "API latency & error rate tracking",
    ]),
    ("Model Refresh Cycle", NAVY, Inches(6.75), [
        "Monthly re-training as new outcomes arrive",
        "Champion-challenger framework (live A/B test)",
        "Automated regression test gate before promotion",
        "Quarterly model audit report for stakeholders",
        "Shadow mode: new model scores without acting",
    ]),
    ("Future Enhancements", AMBER, Inches(9.95), [
        "Bureau pull integration (CIBIL / Experian API)",
        "GST data enrichment for revenue verification",
        "Bank statement PDF auto-parsing (OCR pipeline)",
        "WhatsApp chatbot for borrower self-service scoring",
        "Lender explainability portal (reason codes)",
    ]),
]

for title, bg, lft, items in prod_sections:
    mw2 = Inches(2.9)
    tbar = slide.shapes.add_shape(1, lft, Inches(1.25), mw2, Inches(0.52))
    tbar.fill.solid()
    tbar.fill.fore_color.rgb = bg
    tbar.line.fill.background()
    ttf = tbar.text_frame
    ttf.margin_top = Inches(0.08)
    tp  = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr  = tp.add_run()
    tr.text = title
    tr.font.bold = True
    tr.font.size = Pt(11)
    tr.font.color.rgb = WHITE
    add_bullet_box(slide, lft, Inches(1.8), mw2, Inches(4.8),
                   items, size=9.5)

add_textbox(slide, Inches(0.35), Inches(6.7), Inches(12.6), Inches(0.65),
            "The system is designed to be self-improving: as FinnUp accumulates more outcome data "
            "(disbursements, defaults, rejections), the model iteratively becomes more accurate — "
            "creating a competitive moat for the platform.",
            size=10.5, color=NAVY, bold=True, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — TIMELINE (GANTT)
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Delivery Timeline",
                  "8-week roadmap from today (April 2026) to live scoring in production")

fig, ax = plt.subplots(figsize=(12.5, 5.5), facecolor="white")

phases = [
    # (label, start_week, duration_weeks, color, status)
    ("Phase 1: EDA\n& Data Audit",            0, 2, "#16A34A",  "DONE"),
    ("Phase 2: Feature Engineering",           2, 2, "#F59E0B",  "NEXT"),
    ("Phase 2: Label Creation",                2, 1.5, "#F59E0B","NEXT"),
    ("Phase 3: Baseline Model",                4, 1, "#7C3AED",  "PLANNED"),
    ("Phase 3: XGBoost + Tuning",              4.5, 1.5, "#7C3AED","PLANNED"),
    ("Phase 3: Validation + Model Card",       5.5, 0.8, "#7C3AED","PLANNED"),
    ("Phase 4: REST API Development",          4, 2, "#1B3A6B",  "PLANNED"),
    ("Phase 4: FinnUp Integration / UAT",      6, 1.2, "#1B3A6B","PLANNED"),
    ("Phase 5: Production Deploy",             7, 0.8, "#0D9488", "PLANNED"),
    ("Phase 5: Monitoring Setup",              7.2, 0.8, "#0D9488","PLANNED"),
    ("⚠ Data Labels from FinnUp (BLOCKER)",   1.8, 0.15, "#EF4444","BLOCKER"),
]

yticks = []
yticklabels = []
bar_h = 0.6
gap   = 0.3
for i, (label, start, dur, color, status) in enumerate(phases):
    y = (len(phases) - 1 - i) * (bar_h + gap)
    yticks.append(y + bar_h / 2)
    yticklabels.append(label)
    ax.barh(y, dur, left=start, height=bar_h, color=color, edgecolor="white", linewidth=0.5)
    # status badge
    badge = status
    bx = start + dur + 0.05
    ax.text(bx, y + bar_h / 2, badge, va="center", fontsize=7,
            color=color, fontweight="bold")

ax.set_yticks(yticks)
ax.set_yticklabels(yticklabels, fontsize=8.5)
ax.set_xticks(range(9))
ax.set_xticklabels([f"Wk {i}" if i > 0 else "Now\n(Apr)" for i in range(9)], fontsize=9)
ax.set_xlim(-0.1, 9)
ax.axvline(0, color="#EF4444", linestyle="--", linewidth=1.5, label="Today (Apr 2026)")
ax.axvline(2, color="#F59E0B", linestyle="--", linewidth=0.8, alpha=0.6, label="Labels needed by Wk 2")
ax.set_xlabel("Weeks from now", fontsize=9)
ax.set_title("FinnUp Credit Intelligence — Delivery Gantt", fontsize=11, pad=8)
ax.legend(fontsize=8, loc="lower right")
ax.grid(axis="x", alpha=0.3)
fig.tight_layout()
chart_img = fig_to_bytes(fig)
slide.shapes.add_picture(chart_img, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.9))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — TECHNOLOGY STACK
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Technology Stack",
                  "Python-first, cloud-agnostic, production-grade")

stack = [
    ("Data & Storage", NAVY, [
        "Python 3.11+  ·  Pandas  ·  NumPy  ·  OpenPyXL",
        "Source data: FinnUp Excel / API export",
        "Feature store: Parquet files (→ Postgres in prod)",
        "Model artefacts: MLflow / pickle + version tagging",
    ]),
    ("Machine Learning", TEAL, [
        "scikit-learn (preprocessing, pipelines, cross-val)",
        "XGBoost v2 / LightGBM v4 (tree models)",
        "imbalanced-learn (SMOTE for class imbalance)",
        "SHAP (explainability & feature importance)",
        "Optuna (hyperparameter tuning)",
    ]),
    ("API & Integration", PURPLE, [
        "FastAPI (async REST API, auto OpenAPI docs)",
        "Pydantic v2 (request/response validation)",
        "Celery + Redis (async batch scoring jobs)",
        "Docker (containerisation)",
        "Nginx (reverse proxy)",
    ]),
    ("Cloud & DevOps", AMBER, [
        "AWS / Azure (cloud-agnostic Dockerfile)",
        "GitHub Actions (CI/CD, automated tests)",
        "Terraform / Bicep (infra as code)",
        "Prometheus + Grafana (metrics & alerting)",
        "Sentry (error tracking)",
    ]),
    ("Visualisation & Reports", GREEN, [
        "Jupyter Notebooks (EDA & model analysis)",
        "Matplotlib / Seaborn / Plotly (charts)",
        "python-pptx (auto PowerPoint reports)",
        "WeasyPrint / ReportLab (PDF scorecards)",
        "Streamlit (internal demo dashboard)",
    ]),
]

col_w3 = Inches(2.35)
for i, (title, bg, items) in enumerate(stack):
    cx = Inches(0.35) + i * (col_w3 + Inches(0.15))
    tbar = slide.shapes.add_shape(1, cx, Inches(1.25), col_w3, Inches(0.48))
    tbar.fill.solid()
    tbar.fill.fore_color.rgb = bg
    tbar.line.fill.background()
    ttf = tbar.text_frame
    ttf.margin_top = Inches(0.07)
    tp  = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr  = tp.add_run()
    tr.text = title
    tr.font.bold = True
    tr.font.size = Pt(10)
    tr.font.color.rgb = WHITE
    add_bullet_box(slide, cx, Inches(1.76), col_w3, Inches(5.3),
                   items, size=9.5)

add_textbox(slide, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.3),
            "All components are open-source or self-hostable. No vendor lock-in. FinnUp retains full ownership of models and data.",
            size=9.5, color=GRAY, wrap=False)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — DATA NEEDED FROM FINNUP
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "What We Need from FinnUp",
                  "Blockers and open items required to proceed to modelling")

needed = [
    ("🔴 CRITICAL — BLOCKERS", SALMON, [
        ("Outcome Labels",
         "Approved / Rejected / Disbursed status for each borrower-loan pair. "
         "This is the single most important input for supervised ML training. "
         "Even a partial labelled set of 500–1,000 records would unblock model development.",
         "ASK: Export labelled loan outcomes from CRM / LOS — even partial"),
        ("Lender Policy Rules",
         "Each lender's minimum requirements (CIBIL, turnover, sector, loan size). "
         "These define the eligibility layer that wraps around the ML score.",
         "ASK: Share lender_policies/ folder or Google Sheet with eligibility criteria"),
    ]),
    ("🟡 IMPORTANT — FOR FULL ACCURACY", AMBER, [
        ("Bureau / CIBIL API Access",
         "Direct API to pull real-time CIBIL scores enriches accuracy significantly. "
         "Currently we have only 55 actual borrower CIBIL records — rest are imputed.",
         "ASK: Provide CIBIL / Experian API credentials (or confirm bureau availability)"),
        ("Bank Statement PDFs / AA Consent",
         "Structured bank data currently covers only 38% of borrowers. "
         "Account Aggregator (AA) framework integration or PDF parsing would extend coverage.",
         "ASK: Are borrowers consenting to AA data pull? Share sample PDFs for parser build"),
    ]),
    ("⚪ NICE TO HAVE", TEAL, [
        ("Historical Default / Delinquency Data",
         "Any NPA or delinquency flags in the existing portfolio would strengthen the model.",
         "ASK: Any historical bad loan flags in the data?"),
        ("New Borrower API / Webhook",
         "To enable real-time scoring, FinnUp's tech team needs to expose a webhook on new registration.",
         "ASK: Share API documentation or contact FinnUp engineering team"),
    ]),
]

cur_y = Inches(1.25)
for section_title, bg, items in needed:
    # Section header
    sh_bar = slide.shapes.add_shape(1, Inches(0.35), cur_y, Inches(12.6), Inches(0.40))
    sh_bar.fill.solid()
    sh_bar.fill.fore_color.rgb = bg
    sh_bar.line.fill.background()
    stf = sh_bar.text_frame
    stf.margin_left = Inches(0.12)
    stf.margin_top  = Inches(0.06)
    sp  = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sr  = sp.add_run()
    sr.text = section_title
    sr.font.bold = True
    sr.font.size = Pt(11)
    sr.font.color.rgb = WHITE
    cur_y += Inches(0.43)
    for item_title, item_desc, ask in items:
        add_textbox(slide, Inches(0.5), cur_y, Inches(2.4), Inches(0.55),
                    item_title, bold=True, size=10, color=NAVY)
        add_textbox(slide, Inches(2.95), cur_y, Inches(7.2), Inches(0.55),
                    item_desc, size=9, color=DARK, wrap=True)
        add_textbox(slide, Inches(10.2), cur_y, Inches(3.0), Inches(0.55),
                    ask, size=8.5, color=TEAL, wrap=True)
        cur_y += Inches(0.58)
    cur_y += Inches(0.1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — NEXT STEPS & ACTIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = add_slide(prs, "Next Steps & Immediate Actions",
                  "Prioritised action list — who does what and by when")

actions = [
    # (owner, action, deadline, status_color, priority)
    ("FinnUp Team",
     "Share outcome labels: approved / disbursed / rejected per borrower-loan pair.\n"
     "Even a partial export of 500+ records from the CRM unblocks ML training.",
     "Week 1\n(ASAP)", SALMON, "🔴 P0"),

    ("FinnUp Team",
     "Share lender policy files (already in /data/lender_policies/ folder —\n"
     "confirm if complete or if new lenders need to be added).",
     "Week 1", AMBER, "🔴 P0"),

    ("Data Science Team",
     "Complete feature engineering notebook (02_feature_engineering.ipynb):\n"
     "Engineer 50+ features from all 11 sheets, handle missing values, encode categoricals.",
     "Week 1–2", AMBER, "🟡 P1"),

    ("Data Science Team",
     "Train baseline logistic regression model on FinnUp Status label,\n"
     "establish AUC-ROC benchmark, identify top features via permutation importance.",
     "Week 2–3", TEAL, "🟡 P1"),

    ("Data Science Team",
     "Train XGBoost model with SHAP explanations, hyperparameter tuning via Optuna,\n"
     "temporal train-test split, fairness analysis. Generate model card.",
     "Week 3–4", TEAL, "🟡 P1"),

    ("Engineering",
     "Build FastAPI scoring endpoint (POST /v1/score).\n"
     "Docker containerise. Connect to trained model artefact. Add auth + rate limiting.",
     "Week 5–6", PURPLE, "⚪ P2"),

    ("FinnUp Tech Team",
     "Configure webhook or batch call from FinnUp CRM / portal to scoring API.\n"
     "UAT testing with 100 sample borrowers. Sign off on scorecard format.",
     "Week 6–7", PURPLE, "⚪ P2"),

    ("All Teams",
     "Production deployment, monitoring setup, stakeholder demo.\n"
     "Quarterly model refresh schedule agreed. Go-live.",
     "Week 8", GREEN, "⚪ P3"),
]

hdrs2 = ["Owner", "Action", "Deadline", "Pri"]
col_x3     = [Inches(0.3), Inches(2.3), Inches(10.6), Inches(12.2)]
col_w3_tbl = [Inches(1.9), Inches(8.2), Inches(1.5),  Inches(1.0)]
row_h3     = Inches(0.1)
top3       = Inches(1.28)
hrow_h     = Inches(0.38)

for j, (hdr, xp, wd) in enumerate(zip(hdrs2, col_x3, col_w3_tbl)):
    hb = slide.shapes.add_shape(1, xp, top3, wd, hrow_h)
    hb.fill.solid()
    hb.fill.fore_color.rgb = NAVY
    hb.line.fill.background()
    htf = hb.text_frame
    htf.margin_left = Inches(0.06)
    htf.margin_top  = Inches(0.06)
    hp  = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr  = hp.add_run()
    hr.text = hdr
    hr.font.bold  = True
    hr.font.size  = Pt(10)
    hr.font.color.rgb = WHITE

row_heights = [Inches(0.66), Inches(0.54), Inches(0.66), Inches(0.54),
               Inches(0.66), Inches(0.54), Inches(0.54), Inches(0.54)]
cur_y3 = top3 + hrow_h
for i, ((owner, action, deadline, clr, pri), rh) in enumerate(zip(actions, row_heights)):
    bg = LIGHT if i % 2 == 0 else WHITE
    vals = [owner, action, deadline, pri]
    for j, (val, xp, wd) in enumerate(zip(vals, col_x3, col_w3_tbl)):
        cb = slide.shapes.add_shape(1, xp, cur_y3, wd, rh)
        cb.fill.solid()
        if j == 0:
            cb.fill.fore_color.rgb = clr
        else:
            cb.fill.fore_color.rgb = bg
        cb.line.fill.background()
        ctf = cb.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.06)
        ctf.margin_top  = Inches(0.05)
        cp  = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.LEFT
        cr  = cp.add_run()
        cr.text = val
        cr.font.size  = Pt(9 if j == 1 else 9.5)
        cr.font.bold  = (j == 0 or j == 3)
        cr.font.color.rgb = WHITE if j == 0 else DARK
    cur_y3 += rh


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — CLOSING / THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = NAVY

band2 = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, SH)
band2.fill.solid()
band2.fill.fore_color.rgb = NAVY
band2.line.fill.background()

acc2 = slide.shapes.add_shape(1, Inches(0), Inches(3.5), SW, Pt(5))
acc2.fill.solid()
acc2.fill.fore_color.rgb = TEAL
acc2.line.fill.background()

t_c1 = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(1.0))
cp1  = t_c1.text_frame.paragraphs[0]
cp1.alignment = PP_ALIGN.CENTER
cr1  = cp1.add_run()
cr1.text = "FinnUp — MSME Credit Intelligence"
cr1.font.bold  = True
cr1.font.size  = Pt(32)
cr1.font.color.rgb = WHITE

t_c2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10), Inches(0.6))
cp2  = t_c2.text_frame.paragraphs[0]
cp2.alignment = PP_ALIGN.CENTER
cr2  = cp2.add_run()
cr2.text = "From raw borrower data to automated lending decisions"
cr2.font.size  = Pt(16)
cr2.font.color.rgb = LTBLUE

summary_bullets = [
    "✓  EDA complete — 6,483 borrowers, 11 sheets, 30+ ML-ready features identified",
    "✓  Solution architecture designed — plug-in API / batch pipeline for FinnUp platform",
    "⏳  Waiting: Outcome labels & lender policy confirmation from FinnUp team",
    "◌  8-week delivery plan: Features → Model → API → Production",
]
t_c3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(2.5))
t_c3.text_frame.word_wrap = True
first = True
for b in summary_bullets:
    p = t_c3.text_frame.paragraphs[0] if first else t_c3.text_frame.add_paragraph()
    first = False
    r = p.add_run()
    r.text = b
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFF) if b.startswith("⏳") or b.startswith("◌") else RGBColor(0xBB, 0xF7, 0xD0)

t_c4 = slide.shapes.add_textbox(Inches(1.5), Inches(6.6), Inches(10), Inches(0.6))
cp4  = t_c4.text_frame.paragraphs[0]
cp4.alignment = PP_ALIGN.CENTER
cr4  = cp4.add_run()
cr4.text = "Prepared by Data Science Team  |  April 2026  |  Confidential"
cr4.font.size = Pt(10)
cr4.font.color.rgb = GRAY


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\n✓  Saved: {OUTPUT_PATH.resolve()}")
print(f"   Slides: {len(prs.slides)}")
